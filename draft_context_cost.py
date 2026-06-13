"""
DRAFT — "Bigger context, bigger bill" visualization (TABLE version).

A standalone, throwaway deck with a single slide: what it costs to send a given
amount of context to the model, once, across three Claude tiers. Every cell is
pure arithmetic (tokens x per-token rate) on published rates — so the table is
exact. The real-world multipliers (output tokens, agent loops, caching) live in
the caption, flagged as illustrative rather than dressed up as precise cells.

Reuses the Cato design tokens + drawing helpers from build_presentation.py
WITHOUT touching the main deck (build() is guarded by __main__).

Pricing source: Claude API reference, cached 2026-06-04 —
  Opus 4.8   $5 / 1M input,  $25 / 1M output
  Sonnet 4.6 $3 / 1M input,  $15 / 1M output
  Haiku 4.5  $1 / 1M input,  $5  / 1M output
  No long-context premium on current Opus (1M at standard pricing).
  Prompt caching: cache reads ~0.1x base, cache writes ~1.25x (5-min TTL).

Run:    .venv/bin/python draft_context_cost.py
Output: draft_context_cost.pptx
"""

import os

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from build_presentation import (
    new_prs, content_slide, txb, animate_clicks, add_notes,
    HERE, FONT, MONO, FG, ACCENT, GOLD, MUTED, BOX_FILL, CODE_BG, LINE_DIM,
)

# ── pricing (published rates, $ per input token) ──────────────────────────────
RATES = [("Haiku 4.5", 1.0), ("Sonnet 4.6", 3.0), ("Opus 4.8", 5.0)]  # $ / 1M input
SIZES = [("8K", 8_000), ("32K", 32_000), ("128K", 128_000),
         ("200K", 200_000), ("1M", 1_000_000)]


def _fmt(cost):
    """Readable dollar string with enough precision for sub-cent values."""
    if cost < 0.10:
        return f"${cost:.3f}"
    if cost < 1:
        return f"${cost:.2f}"
    return f"${cost:,.2f}"


def _cell(cell, text, *, size=16, color=FG, bold=False, font=FONT,
          align=PP_ALIGN.LEFT, fill=None):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.18)
    cell.margin_right = Inches(0.18)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def s_context_cost(prs):
    """Cost to send context once, by model — exact arithmetic on published rates."""
    s = content_slide(prs, "LAYER 3 · CONTEXT MANAGEMENT", "Bigger context, bigger bill")

    txb(s, "What it costs to send the context — once, input only",
        Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.45),
        size=18, color=MUTED)

    # ── table ─────────────────────────────────────────────────────────────────
    n_rows = len(SIZES) + 1
    n_cols = len(RATES) + 1
    left, top = Inches(1.15), Inches(2.35)
    width, height = Inches(11.0), Inches(3.15)
    gframe = s.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gframe.table
    table.first_row = False
    table.horz_banding = False

    # column widths: context label wider, three price columns equal
    table.columns[0].width = Inches(3.2)
    for c in range(1, n_cols):
        table.columns[c].width = Inches((11.0 - 3.2) / len(RATES))
    for r in range(n_rows):
        table.rows[r].height = Inches(height.inches / n_rows)

    # header row
    _cell(table.cell(0, 0), "Context sent", size=16, bold=True, color=FG, fill=BOX_FILL)
    for c, (name, rate) in enumerate(RATES, start=1):
        hdr = GOLD if name.startswith("Opus") else ACCENT
        _cell(table.cell(0, c), f"{name}\n${int(rate)} / 1M in",
              size=15, bold=True, color=hdr, align=PP_ALIGN.CENTER, fill=BOX_FILL)

    # body rows
    for r, (slabel, tok) in enumerate(SIZES, start=1):
        row_fill = CODE_BG if r % 2 else None
        _cell(table.cell(r, 0), f"{slabel} tokens", size=16, bold=True,
              color=FG, fill=row_fill)
        for c, (name, rate) in enumerate(RATES, start=1):
            cost = tok * rate / 1_000_000
            col = GOLD if name.startswith("Opus") else FG
            _cell(table.cell(r, c), _fmt(cost), size=16, font=MONO,
                  color=col, align=PP_ALIGN.CENTER, fill=row_fill)

    # ── caption: the real-world multipliers (illustrative, flagged) ────────────
    cap = []
    cap.append(txb(s, "That's a single message. Real usage costs more:",
                   Inches(0.55), Inches(5.75), Inches(12.2), Inches(0.4),
                   size=16, bold=True, color=FG))
    cap.append(txb(s,
        "output tokens 5x input   ·   agents re-send the whole context every step   ·   "
        "prompt caching cuts repeats ~10x",
        Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.4),
        size=14, color=MUTED, align=PP_ALIGN.LEFT))
    cap.append(txb(s,
        "Illustrative: a 200K window over a ~30-step agent runs roughly $15–30 of input on "
        "Opus, not $1 (≈ $4 with caching).",
        Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.4),
        size=13, italic=True, color=GOLD))
    cap.append(txb(s,
        "Rates: Claude API pricing, 2026 — Opus 4.8 $5/$25, Sonnet 4.6 $3/$15, Haiku 4.5 $1/$5 per 1M in/out.  "
        "Current Opus has no long-context premium.",
        Inches(0.55), Inches(7.0), Inches(12.2), Inches(0.34),
        size=11, italic=True, color=MUTED))

    animate_clicks(s, [cap])

    add_notes(s,
        "DRAFT visualization (table) — cost of context.\n\n"
        "Every cell is exact: tokens x the published per-token rate. Reading down a\n"
        "column, cost scales linearly with context; reading across, it scales with the\n"
        "model tier. Sending a full 1M-token prompt to Opus once is ~$5; a small 8K\n"
        "prompt is a few cents. (Current Opus has no long-context premium — older\n"
        "models and Gemini charged ~2x above a threshold; worth a mention.)\n\n"
        "(click) The caveat: this is ONE message, input only. Real cost is higher —\n"
        "output is 5x the input rate, and an agent re-sends its whole context every\n"
        "step, so a long loop multiplies the table value by its step count. Prompt\n"
        "caching reads repeated context at ~1/10 price and is the main lever.\n"
        "The $15-30 / $4 agent figure is an order-of-magnitude illustration, not a\n"
        "measured benchmark — say so if asked. The point: context is a cost lever,\n"
        "not just a quality lever.")


def build():
    prs = new_prs()
    s_context_cost(prs)
    out = os.path.join(HERE, "draft_context_cost.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slide)")
    return out


if __name__ == "__main__":
    build()
