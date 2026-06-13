"""
DRAFT ONLY — standalone build of the proposed closing section ("Arc 1").

Builds a separate deck, `closing_draft.pptx`, so the speaker can read the
proposed 10-slide closing without touching the main deck or its build script.
Reuses the Cato styling/helpers from build_presentation.py.

Choices baked into this draft (easy to change later):
  - CL-7 shows BOTH precedents (ATMs + photography) as compact cards.
  - CL-3 shows the felt-vs-measured BARS.
  - CL-9 keeps ALL THREE self-audit questions.

Animations are intentionally OFF here so every element is visible on a plain
read-through (QuickLook/PowerPoint). The real deck would reveal on click.

Run:    .venv/bin/python build_closing_draft.py
Output: closing_draft.pptx
"""

import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from build_presentation import (
    new_prs, add, ph, del_ph, set_ph, content_slide, image_slide,
    txb, rect, add_notes, asset,
    FG, ACCENT, GOLD, CORAL, SKY, MUTED, BOX_FILL, LINE_DIM, CODE_BG,
    FONT, HEAD,
)

HERE = os.path.dirname(os.path.abspath(__file__))


# ── small local helpers ──────────────────────────────────────────────────────

def center_title(s, text, top=Inches(0.55), size=34, color=FG):
    t = s.shapes.add_textbox(Inches(0.55), top, Inches(12.2), Inches(0.95))
    p = t.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = HEAD
    return t


def bullet(s, marker_text, text, left, top, width, size=24, color=FG, mcolor=ACCENT):
    tb = s.shapes.add_textbox(left, top, width, Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    m = p.add_run(); m.text = marker_text
    m.font.size = Pt(size); m.font.bold = True; m.font.color.rgb = mcolor; m.font.name = FONT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
    return tb


# ── CL-1  Section divider ─────────────────────────────────────────────────────

def cl1_divider(prs):
    s = image_slide(prs)
    t = s.shapes.add_textbox(Inches(0.55), Inches(2.7), Inches(12.2), Inches(1.3))
    p = t.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "So… what about us?"
    r.font.size = Pt(52); r.font.bold = True; r.font.color.rgb = FG; r.font.name = HEAD
    txb(s, "We took the tool apart. Now let's talk about the hand that holds it.",
        Inches(0.55), Inches(4.2), Inches(12.2), Inches(0.7),
        size=22, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "CL-1  |  ~0:30. Hard pivot from tools to people. Slow down. "
                 "'You now understand the three pieces. The last piece is the talk's real subject: you.'")


# ── CL-2  The exciting part ───────────────────────────────────────────────────

def cl2_exciting(prs):
    s = content_slide(prs, "THE EXCITING PART", "This is a great time to be a geek")
    y = Inches(2.1)
    for txt in [
        ("▸   ", "Learn faster — the most patient senior engineer you'll ever sit next to."),
        ("▸   ", "Deliver faster — the boring 70% gets done while you think about the hard part."),
        ("▸   ", "More fun — less drudgery, more of the part you actually love."),
    ]:
        bullet(s, txt[0], txt[1], Inches(0.7), y, Inches(12.0), size=24)
        y += Inches(1.05)
    # data chip
    rect(s, Inches(0.7), Inches(5.7), Inches(11.9), Inches(0.7), fill=BOX_FILL, line=LINE_DIM)
    txb(s, "~73% of developers report more flow and less grunt work.  (GitHub, 2024)",
        Inches(0.9), Inches(5.83), Inches(11.5), Inches(0.5), size=16, color=ACCENT)
    add_notes(s, "CL-2  |  ~1:30. Genuine enthusiasm — this is not a doom talk. AI is a tireless "
                 "tutor, a faster path to shipping, and it clears chores so the fun survives.")


# ── CL-3  The honesty beat (METR) ─────────────────────────────────────────────

def cl3_honesty(prs):
    s = content_slide(prs, "BUT LET'S BE HONEST", "Feeling fast is not the same as being fast")
    # two bars; widths proportional to magnitude (20 vs 19 -> similar), colored differently
    bar_x = Inches(0.9)
    txb(s, "What they FELT", bar_x, Inches(2.25), Inches(4), Inches(0.4), size=16, color=MUTED)
    rect(s, bar_x, Inches(2.7), Inches(6.4), Inches(0.85), fill=GOLD)
    txb(s, "+20% faster", bar_x + Inches(0.2), Inches(2.83), Inches(6), Inches(0.6),
        size=24, bold=True, color=RGBColor(0x18, 0x2D, 0x29))

    txb(s, "What was MEASURED", bar_x, Inches(4.0), Inches(5), Inches(0.4), size=16, color=MUTED)
    rect(s, bar_x, Inches(4.45), Inches(6.1), Inches(0.85), fill=CORAL)
    txb(s, "−19% slower", bar_x + Inches(0.2), Inches(4.58), Inches(6), Inches(0.6),
        size=24, bold=True, color=RGBColor(0x18, 0x2D, 0x29))

    txb(s, "Experienced open-source devs, real tasks, randomized.  (METR, 2025)",
        bar_x, Inches(5.5), Inches(11.5), Inches(0.5), size=15, color=MUTED)
    txb(s, "Knowing the difference is your job. Not the model's.",
        Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.6),
        size=22, bold=True, color=ACCENT)
    add_notes(s, "CL-3  |  ~1:30. Show the felt bar, let them nod, then reveal the measured bar. "
                 "The gap is the point: someone must judge what's actually true and good — that's you.")


# ── CL-4  Moravec ─────────────────────────────────────────────────────────────

def cl4_moravec(prs):
    s = content_slide(prs, "WHY THE GAP", "AI got the difficulty order backwards")
    rect(s, Inches(6.655), Inches(2.1), Pt(1.5), Inches(3.6), fill=LINE_DIM)
    txb(s, "Looks hard. Easy for AI.", Inches(0.55), Inches(2.1), Inches(5.7), Inches(0.5),
        size=18, bold=True, color=SKY)
    for i, t in enumerate(["mental math", "recalling syntax", "trivia",
                           "boilerplate", "summarizing"]):
        txb(s, "•  " + t, Inches(0.7), Inches(2.75) + i * Inches(0.55),
            Inches(5.6), Inches(0.5), size=18, color=FG)
    txb(s, "Looks easy. Hard for AI.", Inches(6.95), Inches(2.1), Inches(5.7), Inches(0.5),
        size=18, bold=True, color=ACCENT)
    for i, t in enumerate(["judgment", "taste", "knowing what matters",
                           "reading the room", "owning the call"]):
        txb(s, "•  " + t, Inches(7.1), Inches(2.75) + i * Inches(0.55),
            Inches(5.6), Inches(0.5), size=18, color=FG)
    txb(s, "What it looks most impressive at is often the least valuable.  (Moravec's paradox)",
        Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.6),
        size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "CL-4  |  ~1:30. Chess and math turned out to be the easy part for machines. The "
                 "toddler stuff — judgment, context, care — is the hard part. Your value is the right column.")


# ── CL-5  Polanyi + take-it-apart callback ────────────────────────────────────

def cl5_polanyi(prs):
    s = content_slide(prs, "THE PART YOU CAN'T HAND OVER", "We know more than we can tell")
    txb(s, "We took AI apart so it would stop feeling like magic.",
        Inches(0.55), Inches(2.3), Inches(12.2), Inches(0.7), size=24, color=FG)
    txb(s, "But the thing that makes you good can't be taken apart, written down, or trained "
           "into a model:", Inches(0.55), Inches(3.15), Inches(12.2), Inches(0.9),
        size=22, color=FG)
    txb(s, "your judgment · your feel for this codebase · your taste",
        Inches(0.55), Inches(4.15), Inches(12.2), Inches(0.7),
        size=24, bold=True, color=SKY, align=PP_ALIGN.CENTER)
    txb(s, "That tacit knowledge is yours.  (Polanyi's paradox)",
        Inches(0.55), Inches(5.7), Inches(12.2), Inches(0.6),
        size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "CL-5  |  ~1:45. Callback to the opening: a geek takes things apart to understand "
                 "them. The twist: some things resist disassembly. You can't write down how you "
                 "know a design smells wrong — and that is exactly what a model can't get.")


# ── CL-6  Value moved ─────────────────────────────────────────────────────────

def cl6_value_moved(prs):
    s = content_slide(prs, "SO THE VALUE MOVED", "The bottleneck is not typing anymore")
    pts = [
        "AI gets you ~70%. The hard 30% (edge cases, security, architecture, "
        "“should we even build this”) is yours.",
        "Writing code got cheap. Reviewing and judging it got expensive.  (PR review time up ~91%)",
        "AI raises the floor. You raise the ceiling.",
    ]
    y = Inches(2.3)
    for p in pts:
        bullet(s, "▸  ", p, Inches(0.7), y, Inches(12.0), size=22)
        y += Inches(1.25)
    add_notes(s, "CL-6  |  ~2:00. Not a threat, a relocation. Everyone can produce code now, so the "
                 "scarce thing is the judgment to know which code should exist and whether it's right.")


# ── CL-7  Precedent + centaur ─────────────────────────────────────────────────

def cl7_precedent(prs):
    s = content_slide(prs, "WE'VE SEEN THIS MOVIE", "The tool changes the job, not your worth")
    # two precedent cards on the left, centaur on the right
    def card(left, top, w, h, head_txt, body_txt, hc=SKY):
        rect(s, left, top, w, h, fill=BOX_FILL, line=LINE_DIM)
        txb(s, head_txt, left + Inches(0.22), top + Inches(0.16), w - Inches(0.44),
            Inches(0.5), size=17, bold=True, color=hc)
        txb(s, body_txt, left + Inches(0.22), top + Inches(0.66), w - Inches(0.44),
            h - Inches(0.8), size=15, color=FG)

    card(Inches(0.55), Inches(2.1), Inches(6.0), Inches(1.55), "ATMs",
         "Were going to end bank tellers. Teller jobs grew, then moved to relationships. (Bessen)")
    card(Inches(0.55), Inches(3.85), Inches(6.0), Inches(1.55), "Photography",
         "Was going to kill painting. It freed it into Impressionism.")
    card(Inches(6.85), Inches(2.1), Inches(5.95), Inches(3.3), "Centaur chess",
         "Human + AI + a good process beats the strongest AI alone. (Kasparov)", hc=ACCENT)

    txb(s, "You + AI beats either one alone. Be the centaur.",
        Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.7),
        size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s, "CL-7  |  ~2:00. Two real precedents, then the lesson Kasparov proved: the winning "
                 "configuration is human + machine + good process. Not beat AI, not be replaced — be the centaur.")


# ── CL-8  What stays human ────────────────────────────────────────────────────

def cl8_stays_human(prs):
    s = content_slide(prs, "WHAT STAYS HUMAN", "Two things never move to the model")
    rect(s, Inches(6.655), Inches(2.1), Pt(1.5), Inches(3.4), fill=LINE_DIM)
    # accountability
    txb(s, "ACCOUNTABILITY", Inches(0.55), Inches(2.1), Inches(5.7), Inches(0.5),
        size=18, bold=True, color=ACCENT)
    txb(s, "“A computer can never be held accountable. That is your job as the human in the loop.”",
        Inches(0.55), Inches(2.7), Inches(5.7), Inches(1.6), size=19, italic=True, color=FG)
    txb(s, "— Simon Willison", Inches(0.55), Inches(4.25), Inches(5.7), Inches(0.4),
        size=15, color=MUTED)
    txb(s, "You own the outcome. That never transfers.",
        Inches(0.55), Inches(4.8), Inches(5.7), Inches(0.7), size=18, color=SKY)
    # taste
    txb(s, "TASTE", Inches(6.95), Inches(2.1), Inches(5.7), Inches(0.5),
        size=18, bold=True, color=ACCENT)
    txb(s, "When anyone can make anything, knowing what is worth making is the rare skill.",
        Inches(6.95), Inches(2.7), Inches(5.7), Inches(1.6), size=20, color=FG)
    txb(s, "— echoing Paul Graham", Inches(6.95), Inches(4.25), Inches(5.7), Inches(0.4),
        size=15, color=MUTED)
    add_notes(s, "CL-8  |  ~1:30. Accountability is dignity, not burden: the model can be wrong, only "
                 "you can be responsible. Taste is the new moat — the rare ones know what's worth generating.")


# ── CL-9  Find your value ─────────────────────────────────────────────────────

def cl9_find_value(prs):
    s = content_slide(prs, "FIND YOUR VALUE", "Do what only you can do")
    txb(s, "Even if AI were better at everything, you still win by spending yourself where your "
           "relative edge is — and delegating the rest.  (Comparative advantage)",
        Inches(0.55), Inches(2.1), Inches(12.2), Inches(0.9), size=18, color=MUTED)
    qs = [
        "What problem here do only I really understand?",
        "What would I refuse to ship, even if it passed every test?",
        "What would I build now that it is 10x cheaper to try?",
    ]
    y = Inches(3.4)
    for q in qs:
        txb(s, q, Inches(0.9), y, Inches(11.8), Inches(0.7), size=24, bold=True, color=ACCENT)
        y += Inches(0.95)
    add_notes(s, "CL-9  |  ~2:00. This is the 'hint towards your own value' the talk promised. Don't "
                 "answer for them — ask the three questions, leave a beat after each. Only they can answer.")


# ── CL-10  Capstone ───────────────────────────────────────────────────────────

def cl10_capstone(prs):
    s = image_slide(prs)
    txb(s, "Give me the right context and I shall move the world.",
        Inches(0.55), Inches(1.9), Inches(12.2), Inches(0.6),
        size=18, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    t = s.shapes.add_textbox(Inches(0.55), Inches(3.0), Inches(12.2), Inches(1.3))
    p = t.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "AI is the lever. You are the fulcrum."
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = HEAD
    txb(s, "It can move the world. You decide which way — and you know what context is even "
           "worth giving it.", Inches(1.0), Inches(4.6), Inches(11.3), Inches(1.0),
        size=22, color=FG, align=PP_ALIGN.CENTER)
    add_notes(s, "CL-10  |  ~1:00. Callback to the tips slide. Flip it onto the human: a lever does "
                 "nothing without a fulcrum, and you are it. Hold, then go to Thank you.")


def build():
    prs = new_prs()
    cl1_divider(prs)
    cl2_exciting(prs)
    cl3_honesty(prs)
    cl4_moravec(prs)
    cl5_polanyi(prs)
    cl6_value_moved(prs)
    cl7_precedent(prs)
    cl8_stays_human(prs)
    cl9_find_value(prs)
    cl10_capstone(prs)
    out = os.path.join(HERE, "closing_draft.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
