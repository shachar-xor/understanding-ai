"""
Build "Understanding AI for Geeks" — re-skinned onto the official Cato template.

Loads `Cato Presentation Template.pptx`, deletes its sample slides, and rebuilds
the talk on Cato's dark (black) layout family so the deck inherits Cato branding
(logo, footer, slide numbers, Aptos fonts, dark-teal theme).

Hybrid approach:
  - Native Cato layouts give background + branding + the title placeholder.
  - Body content is positioned with custom text boxes / shapes (full control,
    matching the original deck), drawn with the Cato colour palette.

Run:    .venv/bin/python build_presentation.py
Output: understanding_ai_for_geeks.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsdecls
from pptx.oxml import parse_xml

HERE     = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "Cato Presentation Template.pptx")
ASSETS   = os.path.join(HERE, "assets")

# ── Cato design tokens ───────────────────────────────────────────────────────
BG       = RGBColor(0x18, 0x2D, 0x29)   # dark teal (inherited from layout bg)
FG       = RGBColor(0xF6, 0xF6, 0xF7)   # off-white text
ACCENT   = RGBColor(0x38, 0xD9, 0xA7)   # mint — primary accent
GOLD     = RGBColor(0xFE, 0xC7, 0x56)
PURPLE   = RGBColor(0x99, 0x55, 0xBA)
CORAL    = RGBColor(0xFF, 0x6F, 0x5B)
SKY      = RGBColor(0x55, 0xCB, 0xF2)
MUTED    = RGBColor(0xB5, 0xCB, 0xCB)   # grey-blue — sub-text
CODE_BG  = RGBColor(0x0E, 0x1A, 0x17)   # very dark teal — code/context blocks
CODE_FG  = RGBColor(0xA8, 0xE6, 0xCF)   # mint-tint code text
LINE_DIM = RGBColor(0x2E, 0x4A, 0x44)   # subtle dividers / outlines
BOX_FILL = RGBColor(0x22, 0x40, 0x3A)   # raised box interior on teal

FONT = "Aptos"
HEAD = "Aptos Display"
MONO = "Consolas"

W = Inches(13.33)
H = Inches(7.5)

# Agent walk-through colour aliases
_DIM  = RGBColor(0x5E, 0x7A, 0x73)   # already-seen context (dim teal)
_NEW  = RGBColor(0x6C, 0xFF, 0xB0)   # tool result / new data (bright green)
_AI   = GOLD                          # previous LLM response (amber/gold)
_HEAD = SKY                           # section markers inside context block


# ═══════════════════════════════════════════════════════════════════════════
#  TEMPLATE / LAYOUT PLUMBING
# ═══════════════════════════════════════════════════════════════════════════

def new_prs():
    """Load the Cato template and strip every sample slide (keep masters/layouts)."""
    prs = Presentation(TEMPLATE)
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)
    return prs


def layout(prs, name):
    """Return the slide layout with the given name (raises if renamed/missing)."""
    for master in prs.slide_masters:
        for lay in master.slide_layouts:
            if lay.name == name:
                return lay
    raise KeyError(f"Layout {name!r} not found in template")


def add(prs, layout_name):
    return prs.slides.add_slide(layout(prs, layout_name))


def ph(slide, idx):
    """Get a placeholder by its idx, or None."""
    for p in slide.placeholders:
        if p.placeholder_format.idx == idx:
            return p
    return None


def del_ph(slide, idx):
    """Remove an unused placeholder so it does not show prompt text."""
    p = ph(slide, idx)
    if p is not None:
        p._element.getparent().remove(p._element)


def set_ph(slide, idx, text, size=None, color=None, bold=None, font=None):
    """Write text into a placeholder, keeping its native (Cato) styling by default."""
    p = ph(slide, idx)
    tf = p.text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            if size is not None:  run.font.size = Pt(size)
            if color is not None: run.font.color.rgb = color
            if bold is not None:  run.font.bold = bold
            if font is not None:  run.font.name = font
    return p


def hide(slide):
    """Mark a slide as hidden (skipped in slideshow)."""
    slide._element.set("show", "0")


# ═══════════════════════════════════════════════════════════════════════════
#  DRAWING HELPERS  (custom-positioned body content)
# ═══════════════════════════════════════════════════════════════════════════

def txb(slide, text, left, top, width, height,
        size=28, bold=False, color=None, align=PP_ALIGN.LEFT,
        font=None, italic=False):
    """Add a single-paragraph text box."""
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE   # keep authored size; avoids PowerPoint reflow gaps
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or FG
    run.font.name = font or FONT
    return tf_box


def txb_runs(slide, segments, left, top, width, height,
             size=28, bold=False, align=PP_ALIGN.LEFT, font=None):
    """Single-paragraph text box built from (text, color) segments, so one word
    can be highlighted in a different color than the rest."""
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color in segments:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or FG
        run.font.name = font or FONT
    return tf_box


def txb_lines(slide, lines, left, top, width, height,
              size=28, bold=False, color=None, align=PP_ALIGN.LEFT,
              font=None, line_spacing=None):
    """Multi-paragraph text box: each item in `lines` is its own paragraph, giving
    a guaranteed line break. (A '\\n' inside a single run is stored as a literal
    newline that renderers treat as whitespace and word-wrap, not a hard break.)"""
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or FG
        run.font.name = font or FONT
    return tf_box


def rect(slide, left, top, width, height, fill=None, line=None, line_w=Pt(1)):
    """Add a rectangle shape."""
    shp = slide.shapes.add_shape(1, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    return shp


def oval(slide, left, top, width, height, fill=None, line=None, line_w=Pt(1)):
    """Add an ellipse/circle shape (mirrors rect())."""
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def code_block(slide, text, left, top, width, height, size=15):
    """Monospaced code block with dark background."""
    rect(slide, left, top, width, height, fill=CODE_BG)
    tf_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.12),
        width - Inches(0.4), height - Inches(0.24))
    tf = tf_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = MONO
    run.font.color.rgb = CODE_FG
    return tf_box


def add_pic_fit(slide, path, left, top, max_w, max_h):
    """Add image preserving aspect ratio within max_w x max_h, centered."""
    from PIL import Image as PILImage
    im = PILImage.open(path)
    iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w = iw * ratio
    h = ih * ratio
    offset_x = (max_w - w) / 2
    offset_y = (max_h - h) / 2
    return slide.shapes.add_picture(path, left + offset_x, top + offset_y, w, h)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def animate_clicks(slide, click_groups):
    """
    Add 'Appear' entrance animations triggered on click.

    click_groups : list of click steps; each step is a list of shape objects
                   revealed together on that click, in order.

    Animated shapes stay hidden until their click, then appear. NOTE: python-pptx
    cannot render animations and QuickLook ignores them — verify playback in
    PowerPoint/Keynote.
    """
    cid = [2]   # id 1 = tmRoot, id 2 = mainSeq (reserved)

    def nid():
        cid[0] += 1
        return cid[0]

    def set_node(spid):
        return (
            '<p:set>'
            '<p:cBhvr>'
            f'<p:cTn id="{nid()}" dur="1" fill="hold">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            '</p:cTn>'
            f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
            '<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>'
            '</p:cBhvr>'
            '<p:to><p:strVal val="visible"/></p:to>'
            '</p:set>'
        )

    def effect_par(spid, node_type):
        return (
            '<p:par>'
            f'<p:cTn id="{nid()}" presetID="1" presetClass="entr" presetSubtype="0" '
            f'fill="hold" grpId="0" nodeType="{node_type}">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{set_node(spid)}</p:childTnLst>'
            '</p:cTn>'
            '</p:par>'
        )

    click_pars = []
    for group in click_groups:
        outer, inner = nid(), nid()
        effects = "".join(
            effect_par(shp.shape_id, "clickEffect" if i == 0 else "withEffect")
            for i, shp in enumerate(group)
        )
        click_pars.append(
            '<p:par>'
            f'<p:cTn id="{outer}" fill="hold">'
            '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            '<p:childTnLst><p:par>'
            f'<p:cTn id="{inner}" fill="hold">'
            '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{effects}</p:childTnLst>'
            '</p:cTn></p:par></p:childTnLst>'
            '</p:cTn>'
            '</p:par>'
        )

    timing = (
        f'<p:timing {nsdecls("p")}>'
        '<p:tnLst><p:par>'
        '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
        '<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(click_pars)}</p:childTnLst>'
        '</p:cTn>'
        '<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        '<p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        '</p:seq></p:childTnLst>'
        '</p:cTn>'
        '</p:par></p:tnLst>'
        '</p:timing>'
    )
    slide._element.append(parse_xml(timing))


def asset(name):
    return os.path.join(ASSETS, name)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE SCAFFOLDS  (native Cato layouts + eyebrow/title)
# ═══════════════════════════════════════════════════════════════════════════

def head(slide, eyebrow, title):
    """Populate the native title placeholder: small mint eyebrow + bold title."""
    t = ph(slide, 0)
    tf = t.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    if eyebrow:
        r0 = p0.add_run()
        r0.text = eyebrow
        r0.font.size = Pt(13)
        r0.font.bold = True
        r0.font.color.rgb = ACCENT
        r0.font.name = FONT
        p1 = tf.add_paragraph()
    else:
        p1 = p0
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = FG
    r1.font.name = HEAD


def content_slide(prs, eyebrow, title):
    """Standard content slide: native 'Text Slide black' + custom body region."""
    s = add(prs, "Text Slide black")
    del_ph(s, 1)            # drop body placeholder; body is custom-positioned
    head(s, eyebrow, title)
    return s


def image_slide(prs):
    """Dark slide with no title — for full-image/mystery slides."""
    s = add(prs, "Text Slide black")
    del_ph(s, 0)
    del_ph(s, 1)
    return s


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    """1 — Title slide (Cover Slide Black)"""
    s = add(prs, "Cover Slide Black")
    set_ph(s, 0, "Understanding AI for Geeks", color=FG)
    del_ph(s, 12)   # no subtitle

    # Three "geek" AI agents (Claude, Copilot, Cursor) as silly tilted stickers
    geeks = [
        ("geek_claude_t.png",   0.75, 4.55, 1.75, -13),
        ("geek_cursot_t.png",   5.30, 5.20, 1.65,  12),
        ("geek_copilot_t.png", 10.35, 4.70, 1.45, -16),
    ]
    for name, l, t, wd, rot in geeks:
        p = asset(name)
        if os.path.exists(p):
            pic = s.shapes.add_picture(p, Inches(l), Inches(t), width=Inches(wd))
            pic.rotation = rot

    add_notes(s,
        "⏱ 0:30 | Running: 0:30\n\n"
        "Let the slide sit for a moment. No need to say much.\n"
        "You can say: 'Let's talk about AI — but not in the way you've heard before.'")


def s02_intro_me(prs):
    """2 — Hi, I'm Shachar (the hook / self-intro)"""
    s = content_slide(prs, "ABOUT ME", "Hi, I'm Shachar")

    # Four bullets — minimal text, revealed one per click. Mint marker + white text.
    bullet_shapes = []
    y = Inches(1.95)
    for text in ["A software developer", "5+ years at Cato",
                 "A Lego & The Office fan", "A geek"]:
        tb = s.shapes.add_textbox(Inches(0.55), y, Inches(6.0), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        m = p.add_run(); m.text = "▸   "
        m.font.size = Pt(26); m.font.bold = True; m.font.color.rgb = ACCENT; m.font.name = FONT
        r = p.add_run(); r.text = text
        r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = FG; r.font.name = FONT
        bullet_shapes.append(tb)
        y += Inches(0.95)

    # Hero photo on the right (AI-enhanced version) — enlarged
    add_pic_fit(s, asset("Cato_leggo_ai_enhanced.png"),
                Inches(6.65), Inches(1.55), Inches(6.55), Inches(5.4))

    # Pop-up punchline: Michael Scott also holding a Dundie (mint-framed inset)
    mx, my, mw, mh = Inches(9.0), Inches(5.5), Inches(4.05), Inches(2.28)
    frame = rect(s, mx - Pt(3), my - Pt(3), mw + Pt(6), mh + Pt(6), fill=ACCENT)
    mscott = s.shapes.add_picture(asset("michael_scott.png"), mx, my, mw, mh)

    # Reveal bullets one per click; the Michael Scott pop-up appears together
    # with the "A Lego & The Office fan" bullet (index 2).
    animate_clicks(s, [
        [bullet_shapes[0]],
        [bullet_shapes[1]],
        [bullet_shapes[2], frame, mscott],
        [bullet_shapes[3]],
    ])

    add_notes(s,
        "⏱ 0:45 | Running: 1:15\n\n"
        "Open warm and light — this is the hook. Click through the four points as you\n"
        "talk; let the photo carry the rest.\n"
        "'Hi, my name is Shachar, and as you can see from my picture here, I'm:\n"
        " a software developer; over 5 years at Cato — there's the 5-years-a-Catonian\n"
        " Lego tower; a Lego and The Office fan — yes, that's my Dundie...'\n"
        " (same click: Michael Scott pops up) '...just like Michael's.'  'And a geek —\n"
        " which is what this talk is about.'\n\n"
        "Click-to-reveal is authored as PowerPoint 'Appear' animations — confirm it\n"
        "plays in PowerPoint/Keynote (it won't show in QuickLook previews).")


def s03_geek_origin(prs):
    """3 — As a geek, I learn by taking things apart (build up / break down / fix)"""
    s = image_slide(prs)

    # Centered title — "love" in mint to echo the heart on the client.
    t = s.shapes.add_textbox(Inches(0.55), Inches(0.5), Inches(12.2), Inches(0.95))
    tf = t.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for txt, col in [("I ", FG), ("love", ACCENT), (" to understand things", FG)]:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = col; r.font.name = HEAD

    # (name, left, top, box_w, box_h) — CatoClient is a bit larger (the focal joke).
    specs = [
        ("lego_pieces.jpg",       Inches(0.45), Inches(1.7), Inches(4.0), Inches(3.6)),
        ("m16_parts_cropped.png", Inches(4.66), Inches(1.7), Inches(4.0), Inches(3.6)),
        ("cato-client.png",       Inches(8.70), Inches(1.5), Inches(4.3), Inches(4.2)),
    ]
    pics = []
    for name, left, top, bw, bh in specs:
        p = asset(name)
        if os.path.exists(p):
            pics.append(add_pic_fit(s, p, left, top, bw, bh))
        else:
            box = rect(s, left, top, bw, bh, fill=BOX_FILL, line=LINE_DIM)
            txb(s, f"[ add {name} ]", left, top + bh / 2, bw, Inches(0.5),
                size=12, color=MUTED, align=PP_ALIGN.CENTER)
            pics.append(box)

    # Little green heart on the CatoClient (the "...but I love it" wink),
    # near its top-right corner as a badge.
    heart = s.shapes.add_textbox(Inches(11.75), Inches(1.55), Inches(0.9), Inches(0.8))
    hp = heart.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run(); hr.text = "💚"; hr.font.size = Pt(34)

    final = txb(s,
        "So let's do the same with AI.",
        Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.8),
        size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Reveal images one per click; the heart pops with the CatoClient; then the line.
    animate_clicks(s, [[pics[0]], [pics[1]], [pics[2], heart], [final]])

    add_notes(s,
        "⏱ ~1:00 | Running: ~2:15\n\n"
        "Continue from 'a geek'. Click through the three images as you talk:\n"
        "'As a geek, I've always understood things by building them up' (click: Lego),\n"
        "'breaking them down' (click: M16 disassembled),\n"
        "'and fixing what's broken' (click: CatoClient... yes, ours too; wink + green heart).\n"
        "'I've basically been doing this my whole life.'\n"
        "'So now let's do the same with AI:' (click: closing line) 'take it apart,\n"
        " and learn to think with it and use it better.'\n\n"
        "Click-to-reveal uses PowerPoint 'Appear' animations; confirm playback in\n"
        "PowerPoint/Keynote (it won't show in QuickLook previews).")


def s04a_done_before_grid(prs):
    """4A — We've done this before (concept-card grid variant; HIDDEN backup)"""
    s = content_slide(prs, "THE INSIGHT", "We've done this before")
    hide(s)   # kept as a backup; timeline variant is the live slide

    cards = [
        ("Object-Oriented",       "encapsulation · inheritance · polymorphism"),
        ("Design Patterns",       "Singleton · Factory · Observer · Strategy"),
        ("SOLID",                 "SRP · OCP · LSP · ISP · DIP"),
        ("Functional",            "pure functions · immutability · composition"),
        ("Cloud / Microservices", "containers · orchestration · IaC"),
        ("Agile / DevOps",        "CI/CD · sprints · retros"),
    ]
    lefts = [Inches(0.45), Inches(4.66), Inches(8.87)]
    tops  = [Inches(1.9), Inches(3.6)]
    col_w, card_h = Inches(4.0), Inches(1.5)

    for i, (name, terms) in enumerate(cards):
        left, top = lefts[i % 3], tops[i // 3]
        rect(s, left, top, col_w, card_h, fill=BOX_FILL, line=LINE_DIM)
        txb(s, name, left + Inches(0.22), top + Inches(0.14),
            col_w - Inches(0.44), Inches(0.5), size=18, bold=True, color=ACCENT)
        txb(s, terms, left + Inches(0.22), top + Inches(0.66),
            col_w - Inches(0.44), Inches(0.75), size=13, color=MUTED)

    payoff = txb(s, "We built the language for software. Now let's build one for AI.",
                 Inches(0.55), Inches(5.55), Inches(12.2), Inches(0.7),
                 size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    animate_clicks(s, [[payoff]])

    add_notes(s,
        "⏱ ~1:30 | Running: ~3:45\n\n"
        "AI is a new technology. But handling new technology is our home turf.\n"
        "Every time a new paradigm arrived, WE coined the words that made it thinkable:\n"
        "object-orientation, design patterns, SOLID, cloud, agile. Thinking of a Singleton\n"
        "or a Factory isn't hard for an experienced developer; SOLID is second nature. Once\n"
        "we coined these terms, juniors and seniors alike could think and communicate them.\n"
        "None of it came from a single vendor; the community built it.\n"
        "(click) 'We built the language for software. Now let's build one for AI.'\n\n"
        "Show more than you name: let the grid speak; you only need to mention a few.")


def s04b_done_before_timeline(prs):
    """4B — We've done this before, as a community (timeline; live slide)"""
    s = content_slide(prs, "", "We've done this before, as a community")

    # Chronologically sorted (community-coined vocabulary, with rough years).
    milestones = [
        ("OOP",             "1970s", "encapsulation · inheritance"),
        ("Design Patterns", "1994",  "Singleton · Factory"),
        ("Agile",           "2001",  "sprints · retros"),
        ("SOLID",           "2004",  "SRP · OCP · DIP"),
        ("Cloud",           "2006",  "containers · microservices"),
        ("AI ?",            "now",   "we write this"),
    ]
    box_w, box_h = Inches(1.82), Inches(1.25)
    step    = Inches(2.12)
    x0      = Inches(0.5)
    box_top = Inches(2.55)

    groups = []
    for i, (name, year, terms) in enumerate(milestones):
        left = x0 + step * i
        last = (i == len(milestones) - 1)
        shapes = []
        if i > 0:  # arrow leading into this milestone, revealed with it
            shapes.append(txb(s, "→",
                left - (step - box_w) - Inches(0.13), box_top + Inches(0.32),
                Inches(0.62), Inches(0.6), size=24, color=MUTED, align=PP_ALIGN.CENTER))
        shapes.append(txb(s, year, left, box_top - Inches(0.5), box_w, Inches(0.4),
            size=14, bold=True, color=ACCENT if last else MUTED, align=PP_ALIGN.CENTER))
        shapes.append(rect(s, left, box_top, box_w, box_h,
            fill=BOX_FILL, line=ACCENT if last else LINE_DIM,
            line_w=Pt(1.75) if last else Pt(0.75)))
        shapes.append(txb(s, name, left, box_top + Inches(0.4), box_w, Inches(0.55),
            size=18, bold=True, color=ACCENT if last else FG, align=PP_ALIGN.CENTER))
        shapes.append(txb(s, terms, left - Inches(0.12), box_top + box_h + Inches(0.1),
            box_w + Inches(0.24), Inches(0.7), size=12, color=MUTED, align=PP_ALIGN.CENTER))
        groups.append(shapes)

    payoff = txb(s, "Every shift, the community coined the words. AI is next.",
                 Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.7),
                 size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    groups[-1].append(payoff)   # payoff lands with the AI milestone

    # Each milestone appears on its own click; the payoff lands with AI.
    animate_clicks(s, groups)

    add_notes(s,
        "⏱ ~1:30 | Running: ~3:45\n\n"
        "AI is new, but handling new technology is home turf for us, as a community.\n"
        "Walk the timeline as you click:\n"
        "OOP (1970s), Design Patterns (1994), Agile (2001), SOLID (2004), Cloud (2006)...\n"
        "each time, WE coined the vocabulary that made the idea thinkable and shareable.\n"
        "A Singleton or a Factory isn't hard for an experienced developer; SOLID is second\n"
        "nature. None of it came from a single vendor; the community built the language.\n"
        "(final click) AI is the next stop, and the words are ours to write:\n"
        "'Every shift, the community coined the words. AI is next.'\n\n"
        "Milestones appear one per click; the payoff lands with the AI box. (PowerPoint\n"
        "'Appear' animations; confirm playback in PowerPoint, not QuickLook.)")


def s05_ai_different(prs):
    """5 — But AI is different: fast and closed (clean typographic)"""
    s = image_slide(prs)

    # Centered title
    t = s.shapes.add_textbox(Inches(0.55), Inches(0.6), Inches(12.2), Inches(0.95))
    tp = t.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = "But AI is different"
    tr.font.size = Pt(34); tr.font.bold = True; tr.font.color.rgb = FG; tr.font.name = HEAD

    # Vertical divider (static)
    rect(s, Inches(6.655), Inches(2.2), Pt(1.5), Inches(3.1), fill=LINE_DIM)

    def panel(cx, icon, word):
        col_w = Inches(5.6)
        left = cx - col_w / 2
        ic = s.shapes.add_textbox(left, Inches(2.45), col_w, Inches(1.4))
        ip = ic.text_frame.paragraphs[0]; ip.alignment = PP_ALIGN.CENTER
        ir = ip.add_run(); ir.text = icon; ir.font.size = Pt(66)
        wd = txb(s, word, left, Inches(4.0), col_w, Inches(1.0),
                 size=46, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        return [ic, wd]

    fast_grp   = panel(Inches(3.43), "⏩", "Fast")
    closed_grp = panel(Inches(9.9),  "🔒", "Closed")

    payoff = txb(s, "So we build it ourselves: as a community, vendor-independent.",
                 Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.9),
                 size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Reveal: Fast, then Closed, then the payoff.
    animate_clicks(s, [fast_grp, closed_grp, [payoff]])

    add_notes(s,
        "⏱ ~1:30 | Running: ~5:15\n\n"
        "The catch: AI denies us the two things that let us build a shared language before.\n"
        "(click) FAST: new models and terms every few weeks. The words never get a chance\n"
        "to settle the way 'Singleton' or 'SOLID' did over years.\n"
        "(click) CLOSED: it's vendor-driven and built for lock-in. OpenAI, Anthropic, Google,\n"
        "Meta all compete; none is incentivized to give us a shared vocabulary.\n"
        "(click) 'So we build it ourselves: as a community, vendor-independent.' We can't wait\n"
        "for it to settle, and nobody hands it to us. That is exactly what the rest of this talk\n"
        "does, with three words.\n\n"
        "Fast / Closed / payoff reveal on click (PowerPoint 'Appear'; confirm in PowerPoint).")


def s04_why_hard(prs):
    """4 — Why AI feels hard (superseded by s05_ai_different; kept for reference)"""
    s = content_slide(prs, "THE PROBLEM", "Why AI feels hard")

    img_top = Inches(1.95)
    img_h   = Inches(3.7)
    col_w   = Inches(5.8)

    def add_half_pic(path, left):
        if os.path.exists(path):
            add_pic_fit(s, path, left, img_top, col_w, img_h)

    add_half_pic(asset("fast.png"),   Inches(0.3))
    add_half_pic(asset("closed.png"), Inches(7.0))

    rect(s, Inches(6.55), img_top, Pt(1), Inches(4.7), fill=LINE_DIM)

    txb(s, "Fast", Inches(0.3), Inches(5.85), Inches(5.8), Inches(0.6),
        size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "Moves faster than we can build shared understanding",
        Inches(0.3), Inches(6.45), Inches(5.8), Inches(0.6),
        size=15, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s, "Closed", Inches(7.0), Inches(5.85), Inches(5.8), Inches(0.6),
        size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "Vendor-driven. Each wants lock-in.\nNo incentive to build a common language.",
        Inches(7.0), Inches(6.45), Inches(5.8), Inches(0.6),
        size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 2:00 | Running: 3:00\n\n"
        "FAST: AI moves faster than our ability to build shared understanding around it.\n"
        "New models drop every few weeks. The terminology shifts constantly.\n\n"
        "CLOSED: OpenAI, Anthropic, Google, Meta — they all compete.\n"
        "Each one wants you locked into their platform, their API, their way of thinking.\n"
        "Nobody is incentivized to create a shared vocabulary.\n"
        "That's why it feels like you're learning from scratch every time.\n\n"
        "These two things together are why AI feels hard — not because it's magic.")


def s05_common_language(prs):
    """5 — We've solved this before"""
    s = content_slide(prs, "THE INSIGHT", "We've solved this before")

    concepts = [
        ("Functional Programming", "Pure functions, immutability, composition"),
        ("Object-Oriented",        "Encapsulation, inheritance, polymorphism"),
        ("Design Patterns",        "Singleton, Factory, Observer, Strategy, Decorator..."),
        ("SOLID Principles",       "Single responsibility, Open/closed, Liskov, "
                                   "Interface segregation, Dependency inversion"),
    ]
    y = Inches(2.0)
    for title, sub in concepts:
        txb(s, title, Inches(0.55), y, Inches(4.5), Inches(0.5),
            size=22, bold=True, color=FG)
        txb(s, sub, Inches(0.55), y + Inches(0.48), Inches(4.5), Inches(0.55),
            size=14, color=MUTED)
        y += Inches(1.2)

    txb(s, '"We built Functional Programming, OOP,\nDesign Patterns, SOLID\n'
           '— none of it was invented by\nMicrosoft or Oracle.\n\n'
           'The community did it.  We can do it again."',
        Inches(5.8), Inches(2.0), Inches(7.0), Inches(3.2),
        size=18, color=FG, italic=True)
    txb(s, "Trivially important  ·  Vendor-independent  ·  Transferable  ·  Community-owned",
        Inches(5.8), Inches(5.4), Inches(7.0), Inches(0.5),
        size=13, color=ACCENT)
    add_notes(s,
        "⏱ 2:30 | Running: 5:30\n\n"
        "We've been here before as a profession.\n"
        "OOP didn't come from Microsoft. Design Patterns didn't come from Oracle.\n"
        "The community built these — a shared language that works regardless of\n"
        "which vendor, language, or framework you use.\n\n"
        "That's what we're doing today: LLM, Agent, Context — our attempt at that\n"
        "vocabulary for AI. Simple, transferable, community-owned.")


def s07_breakdown_rings(prs):
    """7 — Break AI into three EQUAL pieces (segmented ring, built arc by arc)"""
    import math
    s = image_slide(prs)

    # Centered title
    t = s.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(12.2), Inches(0.9))
    tp = t.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run(); tr.text = "We'll break AI into three pieces"
    tr.font.size = Pt(34); tr.font.bold = True; tr.font.color.rgb = FG; tr.font.name = HEAD

    cx, cy = 6.66, 4.1          # center, inches
    side = 5.0                  # square holding the ring arc images
    Lr = 2.55                   # label radius (ring outer radius ~2.0")

    def polar_label(text, ang, size, color, two_line=False, w=2.6, rad=None):
        rad_ = rad if rad is not None else Lr
        x = cx + rad_ * math.cos(math.radians(ang))
        y = cy + rad_ * math.sin(math.radians(ang))
        box = s.shapes.add_textbox(Inches(x - w / 2), Inches(y - 0.45), Inches(w), Inches(0.9))
        tf = box.text_frame; tf.word_wrap = True
        lines = text.split("\n") if two_line else [text]
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            rr = p.add_run(); rr.text = ln
            rr.font.size = Pt(size); rr.font.bold = True
            rr.font.color.rgb = color; rr.font.name = FONT
        return box

    # Three equal arcs (pre-rendered images), placed on the same square so they align.
    def arc(name):
        return s.shapes.add_picture(asset(name),
            Inches(cx - side / 2), Inches(cy - side / 2), Inches(side), Inches(side))
    llm_w     = arc("ring_llm.png")       # top (mint)
    agent_w   = arc("ring_agent.png")     # lower-left (sky)
    context_w = arc("ring_context.png")   # lower-right (gold)

    # "AI" in the (transparent) center
    txb(s, "AI", Inches(cx - 0.9), Inches(cy - 0.45), Inches(1.8), Inches(0.9),
        size=40, bold=True, color=FG, align=PP_ALIGN.CENTER)

    llm_lbl     = polar_label("LLM", 270, 24, ACCENT, w=2.0, rad=2.4)
    agent_lbl   = polar_label("Agent", 150, 24, SKY, w=2.2, rad=2.95)
    context_lbl = polar_label("Context\nManagement", 30, 20, GOLD, two_line=True, w=2.6, rad=2.95)

    # Reveal arc by arc; center + "AI" stay put as the arcs build around them.
    animate_clicks(s, [
        [llm_w, llm_lbl],
        [agent_w, agent_lbl],
        [context_w, context_lbl],
    ])

    add_notes(s,
        "⏱ ~0:40 | Running: ~6:00\n\n"
        "AI breaks down into three equal, fundamental pieces. Light them up one at a time:\n"
        "(click) LLM, the core engine.\n"
        "(click) Agent: an LLM with tools and a loop, so it can take actions.\n"
        "(click) Context Management: what we feed and manage around it.\n"
        "Three equal parts of one whole. We'll take each one in turn over the next sections.\n\n"
        "Arcs reveal one per click (PowerPoint 'Appear'; confirm in PowerPoint).")


def s06_pin_reveal(prs):
    """6 — Pin reveal: pin -> disassembled -> assembled"""
    s = content_slide(prs, "THE METAPHOR", "Know your weapon")

    col1, col2, col3 = Inches(0.3), Inches(4.55), Inches(8.8)
    col_w   = Inches(3.9)
    img_top = Inches(1.95)
    img_h   = Inches(4.1)
    label_t = Inches(6.25)

    add_pic_fit(s, asset("pin_transparent.png"), col1, img_top, col_w, img_h)
    txb(s, "The pin", col1, label_t, col_w, Inches(0.5),
        size=16, color=MUTED, align=PP_ALIGN.CENTER)

    txb(s, "→", Inches(4.1), Inches(3.7), Inches(0.5), Inches(0.6),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "→", Inches(8.35), Inches(3.7), Inches(0.5), Inches(0.6),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)

    for col, img, label in [
        (col2, "m16_parts_cropped.png", "Disassembled"),
        (col3, "m16_assembled.png",     "Assembled"),
    ]:
        p = asset(img)
        if os.path.exists(p):
            add_pic_fit(s, p, col, img_top, col_w, img_h)
        else:
            rect(s, col, img_top, col_w, img_h, fill=BOX_FILL, line=LINE_DIM)
            txb(s, f"[ add {img} ]", col, Inches(3.7), col_w, Inches(0.5),
                size=12, color=MUTED, align=PP_ALIGN.CENTER)
        txb(s, label, col, label_t, col_w, Inches(0.5),
            size=16, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 1:30 | Running: 7:00\n\n"
        "Point at the pin: 'Remember this? A cotter pin — פין פזיל.'\n"
        "'Tiny. Costs nothing. But if it's missing, the M16 won't fire.'\n\n"
        "Disassembled: 'In the IDF, before you fire, you break the weapon down\n"
        "completely — knowing each part gives you confidence and control.'\n\n"
        "Assembled: 'That's what we're doing with AI today. By the end you'll know\n"
        "every part — and none of it will feel like magic.'")


def _roadmap_slide(prs, visible_count, title, notes):
    """Layer stack diagram; visible_count layers active (bottom-up)."""
    ALL_LAYERS = [("LLM", ACCENT), ("Agent", FG), ("Context Management", MUTED)]
    s = content_slide(prs, "THE MAP", title)

    left, width = Inches(1.8), Inches(9.5)
    gap, layer_h = Inches(0.1), Inches(1.5)
    y_bottom = Inches(5.9)

    for i, (name, col) in enumerate(ALL_LAYERS):
        y = y_bottom - i * (layer_h + gap)
        active = (i < visible_count)
        box = rect(s, left, y, width, layer_h,
                   fill=BOX_FILL if active else RGBColor(0x15, 0x21, 0x1E),
                   line=col if active else LINE_DIM,
                   line_w=Pt(1.5) if active else Pt(0.5))
        if active:
            txb(s, name, left + Inches(0.3), y + Inches(0.42),
                width - Inches(0.6), Inches(0.7),
                size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_notes(s, notes)


def s07a_roadmap_llm(prs):
    _roadmap_slide(prs, 1, "Layer 1",
        "⏱ 0:20 | Running: 7:20\n\n"
        "'Let's start with the foundation. One box. LLM.'\n"
        "Everything else sits on top of this.")


def s07b_roadmap_agent(prs):
    _roadmap_slide(prs, 2, "Layer 2",
        "⏱ 0:20 | Running: 7:40\n\n"
        "'Agent sits on top of LLM. You can't have an agent without an LLM inside it.'")


def s07c_roadmap_full(prs):
    _roadmap_slide(prs, 3, "Three layers",
        "⏱ 0:40 | Running: 8:00\n\n"
        "Our roadmap. Three layers, bottom to top.\n\n"
        "LLM: the core unit. Agent: LLM + tools + a loop. "
        "Context Management: everything that controls what the LLM sees.\n\n"
        "By the end, every AI product you've used will make sense — they're all\n"
        "combinations of these three things.")


def _blackbox_slide(prs, show_answer):
    s = content_slide(prs, "LAYER 1 · LLM", "A black box with common human knowledge")

    txb(s, "INPUT", Inches(0.55), Inches(2.0), Inches(2), Inches(0.4),
        size=13, color=MUTED)
    code_block(s, '"Cato Networks is the best SASE company in the ___"',
               Inches(0.55), Inches(2.45), Inches(12.2), Inches(0.65), size=16)
    txb(s, "↓", Inches(0), Inches(3.15), Inches(13.33), Inches(0.5),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)

    rect(s, Inches(5.0), Inches(3.7), Inches(3.3), Inches(1.1),
         fill=RGBColor(0x0B, 0x15, 0x12), line=ACCENT, line_w=Pt(1.5))
    txb(s, "LLM", Inches(5.0), Inches(3.85), Inches(3.3), Inches(0.6),
        size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "↓", Inches(0), Inches(4.85), Inches(13.33), Inches(0.5),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "OUTPUT", Inches(0.55), Inches(5.35), Inches(2), Inches(0.4),
        size=13, color=MUTED)

    if show_answer:
        code_block(s, '"Universe!"',
                   Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.65), size=16)
    else:
        rect(s, Inches(0.55), Inches(5.8), Inches(12.2), Inches(0.65), fill=CODE_BG)

    txb(s, "Trained on most of the text ever written.  Thinks in patterns, not rules.  "
           "Its knowledge was frozen at training time.",
        Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.5),
        size=16, color=MUTED, align=PP_ALIGN.CENTER)
    if show_answer:
        add_notes(s,
            "⏱ 0:30 | Running: 9:00\n\n"
            "'Universe! Not world — universe.' Let the room react.\n"
            "'The model has seen enough Cato marketing to know the answer.\n"
            "It's not thinking — it's completing the most probable next token.'")
    else:
        add_notes(s,
            "⏱ 0:30 | Running: 8:30\n\n"
            "Read the prompt out loud. Pause. 'What would you complete this with?'\n"
            "Let the room answer (most say 'world'). Then click to reveal.")


def s08a_llm_blackbox(prs):
    _blackbox_slide(prs, show_answer=False)


def s08a2_llm_blackbox_reveal(prs):
    _blackbox_slide(prs, show_answer=True)


def s08b_autocomplete(prs):
    """8b — Autocomplete can do logic"""
    s = content_slide(prs, "LAYER 1 · LLM", "Autocomplete can do logic")

    L_W, DIV, R_X, R_W = Inches(4.44), Inches(4.54), Inches(4.74), Inches(8.35)
    rect(s, DIV, Inches(2.0), Pt(1), Inches(4.5), fill=LINE_DIM)

    txb(s, "A 4-year-old", Inches(0.3), Inches(2.0), L_W, Inches(0.45),
        size=15, color=MUTED, italic=True)
    y = Inches(2.6)
    for prompt, ans in [('"The sky is ___"', '"blue"'), ('"Dogs say ___"', '"woof"')]:
        txb(s, prompt, Inches(0.3), y, L_W, Inches(0.45), size=20, color=FG, font=MONO)
        txb(s, f"→  {ans}", Inches(0.3), y + Inches(0.48), L_W, Inches(0.4),
            size=20, color=ACCENT, font=MONO)
        y += Inches(1.2)

    txb(s, "Logic  &  knowledge", R_X, Inches(2.0), R_W, Inches(0.45),
        size=15, color=MUTED, italic=True)
    txb(s, '"All humans are mortal.\nSocrates is human.\nTherefore Socrates is ___"',
        R_X, Inches(2.6), R_W, Inches(1.3), size=18, color=FG, font=MONO)
    txb(s, '→  "mortal"', R_X, Inches(3.95), R_W, Inches(0.4),
        size=18, color=ACCENT, font=MONO)
    txb(s, '"I travel to Sirius at 99%\nthe speed of light and return.\n'
           'Compared to my twin, I aged ___"',
        R_X, Inches(4.55), R_W, Inches(1.3), size=18, color=FG, font=MONO)
    txb(s, '→  "less"', R_X, Inches(5.9), R_W, Inches(0.4),
        size=18, color=ACCENT, font=MONO)

    txb(s, "The mechanism is trivial.  The scale is not.",
        Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5),
        size=17, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 2:00 | Running: 11:00\n\n"
        "Walk left to right; ask the room to complete each before revealing.\n"
        "Left: pure pattern matching. Right — Socrates: pure logic completed from\n"
        "everything ever written. Sirius: special relativity, twin paradox — the LLM\n"
        "completes 'less' from the corpus, not from first-principles reasoning.\n"
        "Same mechanism, two very different levels of abstraction.")


def s08c_context_sensitivity(prs):
    """8c — Same prompt, opposite output"""
    s = content_slide(prs, "LAYER 1 · LLM",
                      "The LLM sees the whole prompt — not just the end")
    code_block(s,
        'Prompt:  "A cat walked into a library and asked for a book about"\n'
        'Output:  "...fish. The librarian smiled and pointed to the nature section."',
        Inches(0.55), Inches(2.0), Inches(12.2), Inches(1.1), size=15)
    code_block(s,
        'Prompt:  "A car walked into a library and asked for a book about"\n'
        'Output:  "...roads and highways. The librarian — confused but helpful — reached for an atlas."',
        Inches(0.55), Inches(3.3), Inches(12.2), Inches(1.1), size=15)
    txb(s, "cat  vs  car  — one letter.",
        Inches(0.55), Inches(4.65), Inches(7), Inches(0.6),
        size=22, bold=True, color=ACCENT)
    txb(s, "The model reads every token.\n"
           "The context window is finite — what you put in it is everything.",
        Inches(0.55), Inches(5.45), Inches(12.2), Inches(0.9),
        size=18, color=MUTED, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 1:30 | Running: 12:00\n\n"
        "Read both prompts — one letter different. The model reads every token,\n"
        "left to right; they all matter. This is why prompt engineering exists: one\n"
        "word early can change behaviour for the whole conversation.\n"
        "Key point: context is a finite budget. What you put in it is everything.")


def s08d_hallucination(prs):
    """8d — Hallucination"""
    s = content_slide(prs, "LAYER 1 · LLM",
                      "Hallucination — autocomplete with no good data")
    code_block(s,
        'Prompt:  "Who won the 1987 regional chess\n'
        '          championship in Tulsa, Oklahoma?"\n\n'
        'Output:  "The 1987 championship was won by\n'
        '          David Harrington, a local teacher who..."',
        Inches(0.55), Inches(2.0), Inches(7.5), Inches(2.4), size=15)
    txb(s, "David Harrington\nprobably doesn't exist.",
        Inches(8.3), Inches(2.1), Inches(4.7), Inches(1.2),
        size=26, bold=True, color=ACCENT)
    txb(s, "Hallucination is not a bug.\nIt's autocomplete applied where\n"
           "there's no good data to complete from.",
        Inches(0.55), Inches(4.8), Inches(12), Inches(1.4),
        size=20, color=FG, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 1:30 | Running: 13:30\n\n"
        "David Harrington almost certainly doesn't exist — a plausible name and story,\n"
        "stated confidently, no warning. Not a bug they forgot to fix: when the training\n"
        "data has no good completion, the model still completes — it has to.\n"
        "Takeaway: confident tone means nothing. Verify specific facts.")


# ═══════════════════════════════════════════════════════════════════════════
#  LLM SECTION  (Piece 1 of AI)  — eyebrow "LLM", colour code = mint (ACCENT)
# ═══════════════════════════════════════════════════════════════════════════

# Shared geometry for the two-column "property" slides
PL_X, PL_W = Inches(0.55), Inches(5.55)     # left column
PDIV_X     = Inches(6.45)                    # centre divider
PR_X, PR_W = Inches(6.95), Inches(5.85)      # right column
PLABEL_Y   = Inches(1.95)
PBODY_Y    = Inches(2.7)
PBOTTOM_Y  = Inches(6.55)


def _divider(s):
    return rect(s, PDIV_X, Inches(2.0), Pt(1), Inches(4.3), fill=LINE_DIM)


def _label(s, text, x, w, color=MUTED):
    return txb(s, text, x, PLABEL_Y, w, Inches(0.5), size=15, color=color, italic=True)


def _ex(s, x, y, w, prompt, answer, psize=19, ph_in=0.5, asize=19, acolor=None):
    """One autocomplete example: prompt line(s) + '→ answer'. Returns [shapes]."""
    p = txb(s, prompt, x, y, w, Inches(ph_in), size=psize, color=FG, font=MONO)
    a = txb(s, "→  " + answer, x, y + Inches(ph_in), w, Inches(0.45),
            size=asize, color=acolor or ACCENT, font=MONO)
    return [p, a]


def _phase_breadcrumb(s, active):
    """Static 'Training → Operation' progress pills, top-right; active is 1 or 2."""
    y = Inches(0.6)
    pw, ph_, gap = Inches(1.7), Inches(0.5), Inches(0.55)
    x = Inches(8.7)
    for label, n in [("Training", 1), ("Operation", 2)]:
        on = (n == active)
        rect(s, x, y, pw, ph_, fill=BOX_FILL if on else None,
             line=ACCENT if on else LINE_DIM, line_w=Pt(1.5) if on else Pt(1))
        txb(s, label, x, y + Inches(0.08), pw, Inches(0.4), size=14, bold=on,
            color=ACCENT if on else MUTED, align=PP_ALIGN.CENTER)
        if n == 1:
            txb(s, "→", x + pw, y, gap, ph_, size=20, color=MUTED, align=PP_ALIGN.CENTER)
        x = x + pw + gap


def s08_llm_opener(prs):
    """LLM section opener — Large Language Model: another ML algorithm, a black box."""
    s = image_slide(prs)
    llm = txb(s, "LLM", Inches(0), Inches(0.95), Inches(13.33), Inches(1.2),
              size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=HEAD)
    full = txb(s, "Large Language Model", Inches(0), Inches(2.25), Inches(13.33), Inches(0.7),
               size=32, bold=True, color=FG, align=PP_ALIGN.CENTER)
    ml = txb(s, "just another machine learning algorithm", Inches(0), Inches(2.98),
             Inches(13.33), Inches(0.5), size=20, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    # The black box: opaque, we won't look inside.
    box_w, box_h = Inches(3.0), Inches(1.45)
    box_x, box_y = Inches((13.33 - 3.0) / 2), Inches(3.95)
    box = rect(s, box_x, box_y, box_w, box_h,
               fill=RGBColor(0x0B, 0x15, 0x12), line=ACCENT, line_w=Pt(1.5))
    box_t = txb(s, "LLM", box_x, box_y + Inches(0.4), box_w, Inches(0.7),
                size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    cap = txb(s, "We'll treat it as a black box.", Inches(0), Inches(5.65),
              Inches(13.33), Inches(0.6), size=22, bold=True, color=FG, align=PP_ALIGN.CENTER)

    animate_clicks(s, [[llm], [full], [ml], [box, box_t, cap]])
    add_notes(s,
        "⏱ 0:40 | Running: ~7:30\n\n"
        "The first piece of AI is the LLM. It just means Large Language Model.\n"
        "Demystify it right away: it is not magic, it is just another machine\n"
        "learning algorithm. And for this whole talk we won't open it up, we'll\n"
        "treat it as a black box: text goes in, text comes out.\n"
        "(Reveal is a PowerPoint 'Appear' animation, confirm in PowerPoint.)")


def s08t_training(prs):
    """Phase 1 — Training (uses llm_training.png)."""
    s = content_slide(prs, "LLM", "Training Stage")
    _phase_breadcrumb(s, active=1)
    pic = s.shapes.add_picture(asset("llm_training.png"),
                               Inches(0.6), Inches(2.15), height=Inches(3.95))
    pic.line.color.rgb = ACCENT
    pic.line.width = Pt(1.5)

    head_t = txb(s, "It trains on text\nhumans have written.", Inches(6.85), Inches(2.25),
                 Inches(6.0), Inches(1.2), size=26, bold=True, color=FG)
    nature = txb(s, "Full of mistakes, opinions,\njokes, memes, even trolling.",
                 Inches(6.85), Inches(3.65), Inches(6.0), Inches(1.1), size=20, color=MUTED)
    frozen = txb(s, "And then it freezes in time.", Inches(6.85), Inches(5.05), Inches(6.0),
                 Inches(0.7), size=26, bold=True, color=ACCENT)

    animate_clicks(s, [[pic, head_t], [nature], [frozen]])
    add_notes(s,
        "⏱ 1:15 | Running: ~8:45\n\n"
        "Phase 1, training, happens once, up front. It is fed an enormous slice of\n"
        "everything people have written, books, code, Reddit, news, Wikipedia, and\n"
        "learns the patterns in it. Two things to say out loud:\n"
        " 1. That text is human, so it carries our mistakes, opinions, jokes, memes,\n"
        "    even deliberate trolling. It is a mirror of the corpus, not an oracle.\n"
        " 2. When training ends, it FREEZES. Its knowledge stops at that moment.\n"
        "Then comes phase 2, operation, which is what we actually use.")


def s08op_operation(prs):
    """Phase 2 — Operation: a black box of autocomplete (Cato joke)."""
    s = content_slide(prs, "LLM", "Operation Stage")
    _phase_breadcrumb(s, active=2)

    txb(s, "It just autocompletes text.", Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.5),
        size=24, bold=True, color=ACCENT)

    txb(s, "INPUT", Inches(0.55), Inches(2.25), Inches(2), Inches(0.4), size=13, color=MUTED)
    code_block(s, '"Cato Networks is the best SASE company in the ___"',
               Inches(0.55), Inches(2.65), Inches(12.2), Inches(0.65), size=16)
    txb(s, "↓", Inches(0), Inches(3.35), Inches(13.33), Inches(0.5),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    rect(s, Inches(5.0), Inches(3.85), Inches(3.3), Inches(1.05),
         fill=RGBColor(0x0B, 0x15, 0x12), line=ACCENT, line_w=Pt(1.5))
    txb(s, "LLM", Inches(5.0), Inches(3.98), Inches(3.3), Inches(0.6),
        size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "↓", Inches(0), Inches(4.95), Inches(13.33), Inches(0.5),
        size=28, color=ACCENT, align=PP_ALIGN.CENTER)
    txb(s, "OUTPUT", Inches(0.55), Inches(5.45), Inches(2), Inches(0.4), size=13, color=MUTED)

    rect(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(0.65), fill=CODE_BG)  # empty panel
    ans = txb(s, '"Universe!"', Inches(0.75), Inches(5.95), Inches(12), Inches(0.5),
              size=18, color=CODE_FG, font=MONO)

    animate_clicks(s, [[ans]])
    add_notes(s,
        "⏱ 0:50 | Running: ~9:35\n\n"
        "Operation is what happens every time you use it. Forget the internals,\n"
        "treat it as a black box of autocomplete. Read the prompt out loud, pause,\n"
        "ask the room to complete it (most say 'world'). Click: 'Universe!'\n"
        "It has seen enough Cato marketing to complete it that way. It is not\n"
        "thinking, it is picking the most probable next words. That is the whole job.")


def s08p1_trivial_logic(prs):
    """Property: trivial pattern, and real logic. Prompts first, answers on click."""
    s = content_slide(prs, "LLM", "It can be trivial, it can reason")
    _divider(s)
    sky_p,  sky_a  = _ex(s, PL_X, Inches(2.7), PL_W, '"The sky is ___"', '"blue"', psize=22, ph_in=0.42, asize=22)
    dog_p,  dog_a  = _ex(s, PL_X, Inches(4.4), PL_W, '"Dogs say ___"', '"woof"', psize=22, ph_in=0.42, asize=22)
    soc_p,  soc_a  = _ex(s, PR_X, Inches(2.7), PR_W,
                         '"All humans are mortal. Socrates is\nhuman. Therefore Socrates is ___"',
                         '"mortal"', psize=18, ph_in=0.9, asize=18)
    twin_p, twin_a = _ex(s, PR_X, Inches(4.6), PR_W,
                         '"I fly near light-speed to a star and\nback. Next to my twin, I aged ___"',
                         '"less"', psize=18, ph_in=0.9, asize=18)

    # Per side: both sentences appear at once, then answers one by one; next click
    # brings up the other side's sentences.
    animate_clicks(s, [
        [sky_p, dog_p],     # left: sentences to complete, together
        [sky_a],            # then answers, one by one
        [dog_a],
        [soc_p, twin_p],    # right: sentences to complete, together
        [soc_a],
        [twin_a],
    ])
    add_notes(s,
        "⏱ 1:30 | Running: ~11:05\n\n"
        "Left (click prompts, let the room answer, click): pure pattern matching, a\n"
        "small child completes these. Right (same): the same autocomplete, now landing\n"
        "logic and physics, because it has read enough philosophy and science text to\n"
        "know how those sentences end. The twin paradox 'less' comes from the corpus,\n"
        "not from reasoning. Say it: the mechanism is trivial, the scale is not.")


def s08p2_knowledge_frozen(prs):
    """Property: holds general knowledge, but frozen in time. Prompts first, answers on click."""
    s = content_slide(prs, "LLM", "It knows a lot, up to a point")
    _divider(s)
    jp_p, jp_a = _ex(s, PL_X, Inches(2.7), PL_W, '"The capital of Japan is ___"', '"Tokyo"', psize=19, ph_in=0.45, asize=19)
    ro_p, ro_a = _ex(s, PL_X, Inches(4.0), PL_W, '"Romeo and ___"', '"Juliet"', psize=19, ph_in=0.45, asize=19)
    wa_p, wa_a = _ex(s, PL_X, Inches(5.3), PL_W,
                     '"Water is hydrogen and ___"', '"oxygen"', psize=19, ph_in=0.45, asize=19)

    ne_p, ne_a = _ex(s, PR_X, Inches(2.7), PR_W,
                     '"The news this morning was ___"', '"(it cannot know)"', psize=19, ph_in=0.45, asize=19)
    md_p, md_a = _ex(s, PR_X, Inches(4.4), PR_W,
                     '"The latest model released is ___"', '"(whatever was current then)"',
                     psize=19, ph_in=0.45, asize=19)

    # Per side: sentences appear together, then answers one by one; next click
    # brings up the other side's sentences.
    animate_clicks(s, [
        [jp_p, ro_p, wa_p],   # left: sentences to complete, together
        [jp_a],               # then answers, one by one
        [ro_a],
        [wa_a],
        [ne_p, md_p],         # right: sentences to complete, together
        [ne_a],
        [md_a],
    ])
    add_notes(s,
        "⏱ 1:15 | Running: ~12:20\n\n"
        "Left (prompts, then answers): a huge amount of settled general knowledge,\n"
        "geography, literature, science, all from the training text. Right: but it is\n"
        "frozen. Ask about today's news or the newest release and it cannot know, it\n"
        "can only guess from whatever was true when training ended (the 'knowledge\n"
        "cutoff'). Say it: it learned the world once, then stopped.")


def s08p3_statistical_predicts(prs):
    """Property: statistical (same prompt, different answers) + predicts, not looks up."""
    s = content_slide(prs, "LLM", "It's not deterministic, it makes guesses")
    _divider(s)

    coffee_p = txb(s, '"A good name for a coffee shop: ___"',
                   PL_X, Inches(2.4), PL_W, Inches(0.45), size=19, color=FG, font=MONO)
    a1 = txb(s, '→  "The Daily Grind"', PL_X, Inches(2.95), PL_W, Inches(0.45),
             size=19, color=ACCENT, font=MONO)
    a2 = txb(s, '→  "Bean There, Done That"', PL_X, Inches(3.55), PL_W, Inches(0.45),
             size=19, color=SKY, font=MONO)

    acct_p, acct_a = _ex(s, PR_X, Inches(2.4), PR_W,
                         '"My account number is ___"', '"(it invents one)"', psize=19, ph_in=0.45, asize=19)
    wx_p, wx_a     = _ex(s, PR_X, Inches(4.0), PR_W,
                         '"The weather right now is ___"', '"(it cannot look)"', psize=19, ph_in=0.45, asize=19)

    # Left side is one sentence with two answers (one per click); next click brings
    # up the right side's sentences together, then their answers one by one.
    animate_clicks(s, [
        [coffee_p],          # left: the sentence to complete
        [a1],                # one answer
        [a2],                # run it again, a different answer
        [acct_p, wx_p],      # right: sentences to complete, together
        [acct_a],
        [wx_a],
    ])
    add_notes(s,
        "⏱ 1:15 | Running: ~13:35\n\n"
        "Left: show the prompt, take a guess from the room, click the first answer,\n"
        "then click again, a different answer for the same prompt. It samples from\n"
        "probabilities, it is not deterministic. Right: that is because it predicts\n"
        "likely text, it does not look anything up. Ask for your account number or\n"
        "the live weather and it will invent or guess. Probability, not retrieval.")


def s08p4_subjective_oneletter(prs):
    """Property: subjective (opinions) + confidently wrong (hallucination)."""
    s = content_slide(prs, "LLM", "It can be subjective, it can be wrong")
    _divider(s)
    band_p, band_a = _ex(s, PL_X, PBODY_Y, PL_W,
                         '"The greatest band ever is ___"', '"The Beatles"', psize=18, asize=18)
    lang_p, lang_a = _ex(s, PL_X, Inches(4.35), PL_W,
                         '"The best programming\nlanguage is ___"', '"Python"',
                         psize=18, ph_in=0.95, asize=18)

    tulsa_p = txb(s, '"Who won the 1987 Tulsa\nchess championship?"',
                  PR_X, PBODY_Y, PR_W, Inches(1.0), size=18, color=FG, font=MONO)
    tulsa_a = txb(s, '→  "David Harrington, a local\n    teacher who..."',
                  PR_X, Inches(3.8), PR_W, Inches(0.95), size=18, color=ACCENT, font=MONO)
    tulsa_note = txb(s, "There was no such championship.\nThe answer is wrong because the\nquestion is wrong.",
                     PR_X, Inches(5.05), PR_W, Inches(1.2), size=17, color=MUTED)

    animate_clicks(s, [
        [band_p, lang_p],   # left: both prompts together
        [band_a],           # left answers one by one
        [lang_a],
        [tulsa_p],          # right prompt
        [tulsa_a],          # the confident (wrong) answer
        [tulsa_note],       # the reveal: the question itself was bogus
    ])
    add_notes(s,
        "⏱ 1:30 | Running: ~15:05\n\n"
        "Left: ask it for 'the best' anything and it answers with an opinion, the\n"
        "majority view of its training text. There is no single truth in there.\n"
        "Right: there was no 1987 Tulsa chess championship. The premise is false,\n"
        "so the confident answer is false too. It does not push back on a bad\n"
        "question, it just completes it. Tone is not correctness.")


def s08p5_every_letter(prs):
    """Property: one letter changes everything (cat vs car), full slide."""
    s = content_slide(prs, "LLM", "Every letter counts")

    CW = Inches(12.2)
    # Top: cat
    cat_p = txb_runs(s, [('"A ', FG), ('cat', ACCENT),
                         (' walked into a library and asked for a book about ___"', FG)],
                     PL_X, Inches(2.25), CW, Inches(0.6),
                     size=22, font=MONO, align=PP_ALIGN.CENTER)
    cat_a = txb(s, '→  "...fish."', PL_X, Inches(3.15), CW, Inches(0.5),
                size=24, bold=True, color=ACCENT, font=MONO, align=PP_ALIGN.CENTER)
    # Bottom: car
    car_p = txb_runs(s, [('"A ', FG), ('car', SKY),
                         (' walked into a library and asked for a book about ___"', FG)],
                     PL_X, Inches(4.7), CW, Inches(0.6),
                     size=22, font=MONO, align=PP_ALIGN.CENTER)
    car_a = txb(s, '→  "...roads and highways."', PL_X, Inches(5.6), CW, Inches(0.5),
                size=24, bold=True, color=SKY, font=MONO, align=PP_ALIGN.CENTER)

    animate_clicks(s, [[cat_p, car_p], [cat_a], [car_a]])
    add_notes(s,
        "⏱ 1:00 | Running: ~16:05\n\n"
        "Same sentence, one letter apart: cat versus car. Read the first one, let\n"
        "the room guess, reveal 'fish'. Then change a single letter and the whole\n"
        "completion changes to 'roads and highways'. It reads every token, and each\n"
        "one steers the guess. This is why prompt wording matters so much: one\n"
        "letter early can change the entire response.")


def s08p6_roleplay(prs):
    """Property: role-play / persona conditioning. Bridges LLM into Agents.

    Both columns are deliberately 'Role + Capabilities + Task': the wizard's
    spells and the engineer's terminal are only flavor text here, the model
    cannot actually use them. That gap is exactly what an agent closes (next
    section), so the speaker narrates the bridge."""
    s = content_slide(prs, "LLM", "It can role-play")
    _divider(s)
    d20 = add_pic_fit(s, asset("d20.png"), Inches(11.25), Inches(0.45),
                      Inches(1.75), Inches(1.75))

    # Left: fantasy (D&D)
    dnd_p = txb(s, '"You are Gandalf, a level 12 wizard.\n'
                   'Spells: Fireball, Detect Magic, Light.\n'
                   'A Balrog blocks the bridge. Your move?"',
                PL_X, Inches(2.55), PL_W, Inches(1.5), size=16, color=FG, font=MONO)
    dnd_a = txb(s, '→  "I slam my staff down: \'You shall not pass!\' and ready Fireball."',
                PL_X, Inches(4.45), PL_W, Inches(1.0), size=16, color=ACCENT, font=MONO)

    # Right: real life (professional system prompt)
    dev_p = txb(s, '"You are a senior software engineer.\n'
                   'You write clean, secure Python and explain\n'
                   'your reasoning. Review this for a bug."',
                PR_X, Inches(2.55), PR_W, Inches(1.5), size=16, color=FG, font=MONO)
    dev_a = txb(s, '→  "First, the bug: this loop is O(n^2). Here is the fix, then why it works."',
                PR_X, Inches(4.45), PR_W, Inches(1.0), size=16, color=ACCENT, font=MONO)

    animate_clicks(s, [[d20], [dnd_p], [dnd_a], [dev_p], [dev_a]])
    add_notes(s,
        "⏱ 1:00 | Running: ~17:05\n\n"
        "Say the title, then click: the d20 drops in. 'The geeky kind of role-play.'\n"
        "Give it a role and it plays along. It is still autocomplete, the role just\n"
        "steers which words are likely next. Left: a wizard with spells. Right: the\n"
        "exact same trick doing a real job, a senior engineer persona, which is how\n"
        "the industry shaped early system prompts. Notice both prompts list a ROLE\n"
        "and a set of CAPABILITIES. But here those spells and that terminal are just\n"
        "words, the model cannot actually cast Fireball or run the code. Give the\n"
        "role REAL tools, and a loop to use them, and you have an agent. That is\n"
        "exactly where we go next.")


# ═══════════════════════════════════════════════════════════════════════════
#  AGENT SECTION  (Piece 2 of AI): from "LLM is limited by context" to the loop
# ═══════════════════════════════════════════════════════════════════════════

def s09a0_context_limit(prs):
    """The LLM is limited by its input size (its context window). That window has
    grown fast over the years, shown as a small chart revealed on a click."""
    s = content_slide(prs, "AGENT", "An LLM is limited by its input size")
    txb_runs(s, [("Its input is the ", MUTED), ("context window", SKY),
                 (": everything it sees at once, measured in tokens.", MUTED)],
             Inches(0.8), Inches(1.95), Inches(11.73), Inches(0.6),
             size=22, align=PP_ALIGN.CENTER)
    cap = txb(s, "Finite, but it keeps growing:  about 2K → 2M tokens in five years.",
              Inches(0.8), Inches(2.85), Inches(11.73), Inches(0.5),
              size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    chart = s.shapes.add_picture(asset("context_growth.png"),
                                 Inches(2.1), Inches(3.25), Inches(9.13), Inches(3.8))
    animate_clicks(s, [[cap, chart]])
    add_notes(s,
        "⏱ 0:50 | Running: ~17:45\n\n"
        "Bridge from role-play: an LLM only ever sees its input, the context window,\n"
        "the system prompt plus your message plus its own answer so far, all as tokens.\n"
        "A fixed budget. Make the point, then click: the small chart shows how that\n"
        "budget has grown, 2K in 2020 to 2M today on a log scale, about a thousand-fold\n"
        "in five years. Steepest jump 2023 to 2024. Obvious thought: just wait, it will\n"
        "hold everything soon. Next slide pushes back.")


def s09a0c_too_big(prs):
    """Words-to-tokens device: even the biggest window cannot hold a whole series."""
    s = content_slide(prs, "AGENT", "But some things still do not fit")

    # The books (left) and the show everyone recognizes (right) frame the math.
    add_pic_fit(s, asset("song_of_ice_and_fire_books.png"), Inches(0.35), Inches(2.45),
                Inches(2.9), Inches(3.85))
    card = add_pic_fit(s, asset("Game_of_Thrones_title_card.jpg"),
                       Inches(9.95), Inches(2.6), Inches(3.3), Inches(3.6))
    rect(s, card.left, card.top, card.width, card.height,
         line=LINE_DIM, line_w=Pt(1))

    cx, cw = Inches(3.4), Inches(6.5)   # centre column between the two images
    book = txb_runs(s, [("A Song of Ice and Fire", GOLD), ("   ·   5 books", MUTED)],
                    cx, Inches(2.5), cw, Inches(0.6), size=23, bold=True,
                    align=PP_ALIGN.CENTER)
    words = txb(s, "≈ 1.77 million words", cx, Inches(3.25), cw, Inches(0.7),
                size=37, bold=True, color=FG, align=PP_ALIGN.CENTER)
    tokens = txb_runs(s, [("→   ", MUTED), ("≈ 2.4 million tokens", ACCENT)],
                      cx, Inches(4.2), cw, Inches(0.75), size=33, bold=True,
                      align=PP_ALIGN.CENTER)
    tnote = txb(s, "(about 1.3 tokens per word)", cx, Inches(4.95), cw, Inches(0.4),
                size=15, color=MUTED, align=PP_ALIGN.CENTER)
    punch = txb(s, "Even the largest LLM can't summarize it in a single pass.",
                cx, Inches(5.6), cw, Inches(1.0), size=24, bold=True, color=GOLD,
                align=PP_ALIGN.CENTER)
    animate_clicks(s, [[book, words], [tokens, tnote], [punch]])
    add_notes(s,
        "⏱ 1:00 | Running: ~18:45\n\n"
        "The books on the left, the title card on the right, everyone knows this one. The\n"
        "cute conversion: words times about 1.3 gives tokens. Game of Thrones, five\n"
        "books, is roughly 2.4M tokens, more than the biggest window on Earth. Click\n"
        "through it, then land the line: even the largest LLM can't summarize it in one\n"
        "pass. (And even when something DOES fit, models get 'lost in the middle' and\n"
        "filling the window is expensive.) So we have to split the work.")


def s09mr_mapreduce(prs):
    """Deterministic 'agent': many LLM calls with plain code routing between them."""
    s = content_slide(prs, "AGENT", "Split it up: many LLMs, code in between")
    pos = (Inches(0), Inches(1.7), Inches(13.33), Inches(5.7))
    s.shapes.add_picture(asset("mapreduce_l1.png"), *pos)          # book (shown first)
    l2 = s.shapes.add_picture(asset("mapreduce_l2.png"), *pos)     # chapters
    l3 = s.shapes.add_picture(asset("mapreduce_l3.png"), *pos)     # MAP: LLMs + summaries
    l4 = s.shapes.add_picture(asset("mapreduce_l4.png"), *pos)     # REDUCE: final summary
    animate_clicks(s, [[l2], [l3], [l4]])
    add_notes(s,
        "⏱ 1:10 | Running: ~19:55\n\n"
        "First trick, no loop yet, just many calls with plain code between them. Start\n"
        "with the whole book. Click: split it into chapters. Click: MAP, summarize each\n"
        "chapter with its own LLM call. Click: REDUCE, one more LLM over the summaries\n"
        "to a final summary. The arrows are ordinary code, deterministic, no AI. Honest\n"
        "caveat: each chapter is summarized blind to the others, so cross-chapter\n"
        "context is lost. Fixed pipeline. What if the steps are not known in advance?")


def s09loop_notools(prs):
    """Feed the LLM's output back into its input: a loop with no tools."""
    s = content_slide(prs, "AGENT", "The agent loop")
    # The concrete prompt you actually write: the exit lives in the instruction.
    prompt = txb(s, '"Write a tweet about our product launch. Critique your draft and '
                    'rewrite until it is punchy, on-brand, and at most 140 characters."',
                 Inches(0.9), Inches(2.0), Inches(11.53), Inches(0.8),
                 size=18, color=ACCENT, align=PP_ALIGN.CENTER, font=MONO)
    # The loop itself (image is 16x4.1, so height = 13.33 / 3.902).
    loop = s.shapes.add_picture(asset("agent_loop_notools.png"),
                                Inches(0), Inches(3.05), Inches(13.33), Inches(3.42))
    animate_clicks(s, [[prompt], [loop]])
    add_notes(s,
        "⏱ 0:55 | Running: ~21:20\n\n"
        "A loop with no tools: the model drafts a tweet, critiques its own draft,\n"
        "rewrites, and repeats until it judges the result good. Click for the prompt,\n"
        "click for the loop, the draft feeds straight back into the LLM. The exit lives\n"
        "in the prompt itself (punchy, on-brand, at most 140 characters). Honest catch,\n"
        "and the bridge to the next slide: with no tools it is grading its own homework.\n"
        "It is fine judging 'punchy' and 'on-brand', but it cannot reliably count to 140\n"
        "characters (it sees tokens, not letters), so it may stop while still over the\n"
        "limit. That is exactly what a tool fixes: run code that actually counts.")


def s09loop_tools(prs):
    """Add real tools to the loop and you have an agent (reuses the full diagram)."""
    s = content_slide(prs, "AGENT", "Agent  =  LOOP( LLM + Tools )")
    s.shapes.add_picture(asset("agent_diagram_4.png"),
                         Inches(0), Inches(1.55), Inches(13.33), Inches(4.85))
    tools = txb(s, "Tools:   web search  ·  calculator  ·  shell commands",
                Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5),
                size=18, color=SKY, align=PP_ALIGN.CENTER)
    animate_clicks(s, [[tools]])
    add_notes(s,
        "⏱ 0:45 | Running: ~22:10\n\n"
        "Same loop, but now the model can call real tools: search, calculator, shell\n"
        "commands. Click for the tool list. One question enters at the LLM; at the\n"
        "bottom it asks 'done?'. If no, it calls a tool, reads the result, loops. If\n"
        "yes, it exits as the agent output. LLM plus tools plus a loop. That is an\n"
        "agent.")


def s09cloudlocal(prs):
    """The loop logic runs cloud or local; the model is usually a cloud API."""
    s = content_slide(prs, "AGENT",
                       "The loop runs anywhere. The brain is usually in the cloud.")
    _divider(s)
    cl_label = txb(s, "CLOUD AGENTS", PL_X, PLABEL_Y, PL_W, Inches(0.5),
                   size=18, bold=True, color=SKY)
    cl_list = txb_lines(s, ["ChatGPT agent", "Gemini in Google Search", "Slackbot"],
                        PL_X, Inches(2.75), PL_W, Inches(2.7),
                        size=24, color=FG, line_spacing=1.25)
    lo_label = txb(s, "LOCAL AGENTS", PR_X, PLABEL_Y, PR_W, Inches(0.5),
                   size=18, bold=True, color=ACCENT)
    lo_list = txb_lines(s, ["Claude Code", "Cursor", "Codex CLI"],
                        PR_X, Inches(2.75), PR_W, Inches(2.7),
                        size=24, color=FG, line_spacing=1.25)
    animate_clicks(s, [[cl_label, cl_list], [lo_label, lo_list]])
    add_notes(s,
        "⏱ 0:45 | Running: ~23:00\n\n"
        "The loop, the code between the LLM calls, can run anywhere. Cloud agents run\n"
        "it on someone else's servers (ChatGPT agent, Gemini in Google Search,\n"
        "Slackbot). Local agents run it on your machine (Claude Code, Cursor, Codex\n"
        "CLI). Click each side. Worth saying out loud: either way the model itself is\n"
        "usually a cloud API. We are a local-agent crowd here.")


def _agent_diagram_slide(prs, eyebrow, title, img_name, notes):
    """Full-bleed agent diagram PNG below the title."""
    s = content_slide(prs, eyebrow, title)
    s.shapes.add_picture(asset(img_name),
                         Inches(0), Inches(1.85), Inches(13.33), Inches(5.5))
    add_notes(s, notes)


def s09a1_agent_llm(prs):
    _agent_diagram_slide(prs, "LAYER 2 · AGENT", "Agent  =  LLM",
        "agent_diagram_1.png",
        "⏱ 0:20 | Running: 14:20\n\n"
        "'At its core, an agent starts with just an LLM. It reads input and decides\n"
        "what to do — in text.'")


def s09a2_agent_llm_tools(prs):
    _agent_diagram_slide(prs, "LAYER 2 · AGENT", "Agent  =  LLM  +  Tools",
        "agent_diagram_2.png",
        "⏱ 0:20 | Running: 14:40\n\n"
        "'Add tools — web search, calculator, APIs. The LLM decides which tool to\n"
        "call; the framework runs it.'")


def s09a3_agent_loop(prs):
    _agent_diagram_slide(prs, "LAYER 2 · AGENT", "Agent  =  LOOP( LLM + Tools )",
        "agent_diagram_3.png",
        "⏱ 0:30 | Running: 15:10\n\n"
        "'Wrap it in a loop. The tool result goes back to the LLM, which re-reads\n"
        "everything and decides the next step — until it decides it's done.'")


def s09a4_agent_full(prs):
    _agent_diagram_slide(prs, "LAYER 2 · AGENT", "Agent  =  LOOP( LLM + Tools )",
        "agent_diagram_4.png",
        "⏱ 1:00 | Running: 16:00\n\n"
        "'User input enters at the LLM. At the bottom it asks: done? If yes — output\n"
        "exits as the Agent Output. If no — it loops, calls a tool, reads the result,\n"
        "asks again. One question in; everything else happens inside the loop.'")


# Reusable agent walk-through layout ------------------------------------------

def _agent_slide(prs, call_num, total, context_lines, new_lines, decision,
                 title="Context accumulates with every step",
                 show_decision=True, show_legend=True):
    s = content_slide(prs, f"AGENT  ·  CALL {call_num} of {total}", title)

    box_left, box_top = Inches(0.45), Inches(1.95)
    box_w, box_h = Inches(8.2), Inches(5.25)
    rect(s, box_left, box_top, box_w, box_h, fill=CODE_BG)

    tf_box = s.shapes.add_textbox(
        box_left + Inches(0.18), box_top + Inches(0.12),
        box_w - Inches(0.36), box_h - Inches(0.24))
    tf = tf_box.text_frame
    tf.word_wrap = False
    first = True
    for text, color in context_lines + new_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.name = MONO
        run.font.color.rgb = color

    d_left, d_top, d_w = Inches(8.9), Inches(1.95), Inches(4.1)
    if show_decision:
        txb(s, "LLM decided:", d_left, d_top, d_w, Inches(0.4), size=13, color=MUTED)
        rect(s, d_left, d_top + Inches(0.42), d_w, Inches(1.5),
             fill=RGBColor(0x12, 0x22, 0x1E), line=ACCENT, line_w=Pt(1))
        txb(s, decision, d_left + Inches(0.15), d_top + Inches(0.55),
            d_w - Inches(0.3), Inches(1.2), size=15, color=ACCENT, font=MONO)

    if show_legend:
        txb(s, "▌ tool result (new)", d_left, d_top + Inches(2.2), d_w, Inches(0.4),
            size=13, color=_NEW)
        txb(s, "▌ previous LLM response", d_left, d_top + Inches(2.65), d_w, Inches(0.4),
            size=13, color=_AI)
    return s


def s09b_developer_wrote(prs):
    """9b: open the walk-through with the user's prompt."""
    s = content_slide(prs, "AGENT", "Let's see how an AI agent works")
    label = txb(s, "User prompt:", Inches(2.5), Inches(2.35), Inches(8.3), Inches(0.4),
                size=15, color=MUTED)
    box = rect(s, Inches(2.5), Inches(2.85), Inches(8.3), Inches(1.2),
               fill=CODE_BG, line=_HEAD, line_w=Pt(1))
    prompt = txb(s, "Who is Taylor Swift's boyfriend?\n"
                    "Find out his age and calculate the square root of his age.",
                 Inches(2.7), Inches(3.0), Inches(7.9), Inches(0.85),
                 size=15, color=FG, font=MONO)
    animate_clicks(s, [[label, box, prompt]])
    add_notes(s,
        "⏱ 0:20 | Running: 15:20\n\n"
        "Click to reveal the prompt. (I asked this because my niece is a huge Swiftie.)\n"
        "One plain question. Next we watch the agent work through it call by call.")


_SYS_HUMAN = [
    ("[SYSTEM]",                                                   _HEAD),
    ("You are a helpful research assistant.",                      _NEW),
    ("When given a task:",                                         _NEW),
    ("  1. Search the web for any facts you need.",                _NEW),
    ("  2. Use the calculator for any math.",                      _NEW),
    ("  3. Always show your reasoning step by step.",              _NEW),
    ("",                                                           _NEW),
    ("Available tools:",                                           _NEW),
    ("  web_search(query)    : search the internet",               _NEW),
    ("  calculator(expr)     : evaluate a math expression",        _NEW),
    ("",                                                           _NEW),
    ("[HUMAN]",                                                    _HEAD),
    ("Who is Taylor Swift's boyfriend?",                           _NEW),
    ("Find out his age and calculate the square root of his age.", _NEW),
]


def s09b1_call0(prs):
    """Call 0: the assembled prompt the agent sends on the first call."""
    s = _agent_slide(prs, 0, 4, [], _SYS_HUMAN, "",
                     title="The agent wraps your prompt in context",
                     show_decision=False, show_legend=False)
    txb(s, "This is what the LLM reads on call 1.",
        Inches(8.9), Inches(2.1), Inches(4.1), Inches(1.0),
        size=14, color=MUTED, italic=True)
    add_notes(s,
        "⏱ 0:40 | Running: 16:00\n\n"
        "The agent didn't change the question, it wrapped it: system instructions, a\n"
        "tools list, then the original question nested inside. Point at the tools list:\n"
        "this is how the LLM knows what tools exist, just text, a description, not a\n"
        "function registry. Nothing has run yet. Click to watch call 1.")


def s09c_call1(prs):
    s = _agent_slide(prs, 1, 4, [], _SYS_HUMAN,
                     'web_search(\n  "Taylor Swift\n   boyfriend")',
                     show_legend=False)
    add_notes(s,
        "⏱ 1:00 | Running: 16:45\n\n"
        "Everything is green, first call, all new. The LLM reads top to bottom:\n"
        "I'm a research assistant, I have these tools, here's the question. It decides\n"
        "it needs to search, and outputs a tool call, in text.")


def s09d_call2(prs):
    ctx = [
        ("[SYSTEM]",                                                   _HEAD),
        ("You are a helpful research assistant. ...",                  _DIM),
        ("Available tools: web_search, calculator",                   _DIM),
        ("",                                                           _DIM),
        ("[HUMAN]",                                                    _HEAD),
        ("Who is Taylor Swift's boyfriend? ...",                       _DIM),
        ("",                                                           _DIM),
        ("[LLM]",                                                      _HEAD),
        ("I'll help find information...",                              _DIM),
        ('→ web_search("Taylor Swift boyfriend")',                _AI),
    ]
    new = [
        ("",                                                           _NEW),
        ("[TOOL RESULT]",                                              _HEAD),
        ("Taylor Swift and Travis Kelce confirmed their romance",      _NEW),
        ("in 2023. The couple announced their engagement on",          _NEW),
        ("August 26, 2025...",                                         _NEW),
    ]
    s = _agent_slide(prs, 2, 4, ctx, new,
                     'web_search(\n  "Travis Kelce\n   age born")')
    add_notes(s,
        "⏱ 1:00 | Running: 18:00\n\n"
        "Grey = already seen, amber = previous LLM response, green = new tool result.\n"
        "The search result is appended; the entire context is sent again. The LLM is\n"
        "never told it's on step 2, it figures that out by re-reading from the top.")


def s09e_call3(prs):
    ctx = [
        ("[SYSTEM]  ...",                                              _DIM),
        ("[HUMAN]   Who is Taylor Swift's boyfriend? ...",             _DIM),
        ("[LLM]     I'll help... -> web_search(\"Taylor Swift...\")", _DIM),
        ("[TOOL]    Taylor Swift and Travis Kelce...",                 _DIM),
        ('[LLM]     → web_search("Travis Kelce age born")',       _AI),
    ]
    new = [
        ("",                                                           _NEW),
        ("[TOOL RESULT]",                                              _HEAD),
        ("Born: October 5, 1989 (age 35).",                           _NEW),
        ("Travis Kelce, born in Westlake, Ohio,",                      _NEW),
        ("is currently 35 years old...",                              _NEW),
    ]
    s = _agent_slide(prs, 3, 4, ctx, new, 'calculator(\n  "sqrt(35)")')
    add_notes(s,
        "⏱ 1:00 | Running: 18:45\n\n"
        "Context keeps growing; everything from calls 1-2 is still there and re-read.\n"
        "Now it knows: boyfriend is Kelce, age 35. What's left? The square root,\n"
        "it switches tools to the calculator. Still just text output.")


def s09f_call4(prs):
    ctx = [
        ("[SYSTEM]  ...",                                              _DIM),
        ("[HUMAN]   Who is Taylor Swift's boyfriend? ...",             _DIM),
        ('[LLM]     → web_search("Taylor Swift boyfriend")',      _DIM),
        ("[TOOL]    Taylor Swift and Travis Kelce...",                 _DIM),
        ('[LLM]     → web_search("Travis Kelce age born")',       _DIM),
        ("[TOOL]    Born: October 5, 1989 (age 35)...",                _DIM),
        ("[LLM]     Now let me calculate...",                          _DIM),
        ('[LLM]     → calculator("sqrt(35)")',                    _AI),
    ]
    new = [
        ("",                                                           _NEW),
        ("[TOOL RESULT]",                                              _HEAD),
        ("5.916079783099616",                                          _NEW),
        ("",                                                           _NEW),
        ("[LLM: FINAL ANSWER]",                                        _HEAD),
        ("Taylor Swift's boyfriend is Travis Kelce.",                  _AI),
        ("His age: 35.  √35 ≈ 5.92",                         _AI),
    ]
    s = _agent_slide(prs, 4, 4, ctx, new, "Done.\nNo more tools\nneeded.")
    add_notes(s,
        "⏱ 1:00 | Running: 19:45\n\n"
        "Full history dimmed; new green: calculator result and final answer. The LLM\n"
        "reads four calls' worth of context, sees the result, and decides it has\n"
        "everything, no more tools. One question from the user; 4 LLM calls,\n"
        "2 web searches, 1 calculator call underneath. That's an agent.")


def s10a_context_why(prs):
    """Layer 3 opener: why managing context is the top layer (cost AND quality),
    and the two things to always watch (how full, and what is inside)."""
    s = content_slide(prs, "CONTEXT MANAGEMENT",
                      "Context is king")

    why = txb_runs(s, [("More context  =  ", FG), ("more cost", GOLD), (",  ", FG),
                       ("more latency", FG), (",  and  ", FG), ("worse answers", GOLD),
                       (".", FG)],
                   Inches(0.55), Inches(2.6), Inches(12.2), Inches(0.6),
                   size=26, bold=True, align=PP_ALIGN.CENTER)

    # The "what is inside + how full" picture: one labeled context-window bar.
    cap = txb(s, "We should be aware of the context size, and the context content.",
              Inches(0.55), Inches(4.2), Inches(12.2), Inches(0.6),
              size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    bar_x, bar_y, bar_h = Inches(1.0), Inches(5.05), Inches(0.72)
    segs = [("system + tools", 1.4, SKY),
            ("history", 3.1, ACCENT),
            ("tool results", 3.0, GOLD),
            ("your prompts", 1.1, FG),
            ("free", 2.7, LINE_DIM)]
    bar_shapes = [cap]
    x = bar_x
    for label, w_in, col in segs:
        w = Inches(w_in)
        is_free = (label == "free")
        seg = rect(s, x, bar_y, w, bar_h,
                   fill=BOX_FILL if is_free else col,
                   line=col, line_w=Pt(1.5))
        lbl = txb(s, label, x, bar_y + bar_h + Inches(0.08), w, Inches(0.4),
                  size=13, color=MUTED if is_free else col, align=PP_ALIGN.CENTER)
        bar_shapes += [seg, lbl]
        x += w

    animate_clicks(s, [[why], bar_shapes])
    add_notes(s,
        "⏱ 1:15 | Running: ~23:30\n\n"
        "Callback to the walk-through: every call appended to the context. Why care?\n"
        "More context is not free: it costs more money, adds latency, AND makes answers\n"
        "worse (the model loses the thread in a bloated window). So the job is to manage\n"
        "it. Click the bar: be aware of two things, how full the window is, and what is\n"
        "actually inside it (system, history, tool results, your prompts). That awareness\n"
        "is the whole layer. Next: what to do about it.")


def s10a2_context_moves(prs):
    """The three concrete context-management moves: start fresh, compress, offload."""
    s = content_slide(prs, "CONTEXT MANAGEMENT",
                      "The basics")

    cols = [
        ("Start fresh", ACCENT,
         "the task changed, or the context\nis full of stale junk",
         "clear it, keep only\nwhat still matters"),
        ("Offload to a file", GOLD,
         "you will need it later,\nbut not right now",
         "write it to disk, keep a\npointer, load on demand"),
        ("Compress", SKY,
         "history is long but\nstill relevant",
         "summarize it into a\nshort note (compaction)"),
    ]
    xs = [Inches(0.7), Inches(4.75), Inches(8.8)]
    col_w = Inches(3.8)
    # dividers between the three columns
    rect(s, Inches(4.6), Inches(2.35), Pt(1), Inches(4.0), fill=LINE_DIM)
    rect(s, Inches(8.65), Inches(2.35), Pt(1), Inches(4.0), fill=LINE_DIM)

    groups = []
    for x, (title, col, when, how) in zip(xs, cols):
        head = txb(s, title, x, Inches(2.45), col_w, Inches(0.8),
                   size=27, bold=True, color=col)
        when_b = txb(s, when, x, Inches(3.45), col_w, Inches(1.3),
                     size=19, color=MUTED, italic=True)
        how_b = txb_runs(s, [("→  ", col), (how, FG)],
                         x, Inches(4.85), col_w, Inches(1.3), size=20)
        groups.append([head, when_b, how_b])

    animate_clicks(s, groups)
    add_notes(s,
        "⏱ 1:00 | Running: ~24:30\n\n"
        "Three moves when the window fills. Start fresh: the task changed or the context\n"
        "is full of stale junk, so clear it and keep only what matters. Offload to a file:\n"
        "you will need it later but not now, so write it to disk, keep a pointer, load it\n"
        "back on demand. Compress: history is long but still relevant, summarize it into a\n"
        "short note (this is compaction). The cleanest version of offload-to-a-file is a\n"
        "skill, which is next.")


def s10b_skills_filesystem(prs):
    """Skills: the worked example of offload-to-a-file / load-on-demand."""
    s = content_slide(prs, "CONTEXT MANAGEMENT", "Skills")
    txb(s, "A skill is a markdown file: a short description, then the content.",
        Inches(0.55), Inches(2.15), Inches(12.2), Inches(0.6),
        size=26, bold=True, color=FG, align=PP_ALIGN.CENTER)
    code_block(s,
        'name: commit\n'
        'description: Use when the user wants to make a git commit\n'
        '---\n'
        '1. Stage the relevant changes with git add\n'
        '2. Write a short, clear commit message\n'
        '3. Run the pre-commit hooks, fix any failures\n'
        '4. Commit, then show the result with git log',
        Inches(1.4), Inches(3.15), Inches(10.53), Inches(2.85), size=20)
    txb(s, "description stays in context  ·  body loads only when needed",
        Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.5),
        size=18, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 1:15 | Running: ~25:45\n\n"
        "The cleanest offload-to-a-file: a skill is literally a markdown file, a short\n"
        "description plus a body. Only the description sits in context; the agent\n"
        "pattern-matches it against the task and loads the body on demand (progressive\n"
        "disclosure). Four wins: load/unload (pay for context only when needed), share\n"
        "(commit the folder), community (use others' skills, like npm/pip), vendor-\n"
        "independent (just text). The filesystem is your context-management layer.")


def s10c_ecosystem(prs):
    """10c: the ecosystem, many names for context management."""
    s = content_slide(prs, "CONTEXT MANAGEMENT",
                      "Every vendor has its own tools")
    txb(s, "Different names, same job: deciding what goes in the window.",
        Inches(0.55), Inches(2.0), Inches(12), Inches(0.5), size=20, color=MUTED)
    y = Inches(2.65)
    for name, desc in [
        ("Claude Code",    "/compact to compress, CLAUDE.md as persistent memory"),
        ("Cursor",         ".cursor/rules + codebase index, loaded into each prompt"),
        ("ChatGPT",        "Memory and Projects: saved facts, per-project workspace"),
        ("GitHub Copilot", "copilot-instructions.md: repo context for the agent"),
        ("Gemini",         "Huge context window, plus context caching to reuse it"),
    ]:
        txb(s, name, Inches(0.55), y, Inches(4.2), Inches(0.48),
            size=17, bold=True, color=ACCENT)
        txb(s, desc, Inches(4.9), y + Inches(0.04), Inches(7.9), Inches(0.44),
            size=16, color=FG)
        y += Inches(0.62)
    txb(s, "Each vendor ships its own utilities, with different names and polish.",
        Inches(0.55), Inches(6.0), Inches(12), Inches(0.5),
        size=17, color=MUTED, align=PP_ALIGN.CENTER)
    txb(s, "Understanding and managing the context stays on us.",
        Inches(0.55), Inches(6.55), Inches(12), Inches(0.5),
        size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_notes(s,
        "⏱ 1:30 | Running: 25:30\n\n"
        "Every vendor wraps context management in its own tools: Claude Code has /compact\n"
        "and CLAUDE.md, Cursor has rules plus a codebase index, ChatGPT has Memory and\n"
        "Projects, Copilot has its instructions file, Gemini leans on a huge window plus\n"
        "context caching. Different names and polish, one job: what goes in the window.\n"
        "But the tools only assist, the responsibility to understand and manage the\n"
        "context is on us.")


def s10r_so_this_is_ai(prs):
    """Recap: the three-piece ring (smaller, left) -> "but, where does it leave us?" """
    import math
    s = content_slide(prs, "THE WHOLE PICTURE", "So this is AI")

    # Ring graphic from slide 7, shrunk and shifted to the left half.
    cx, cy = 3.55, 4.45         # center, inches
    side = 4.0                  # square holding the ring arc images

    def polar_label(text, ang, size, color, two_line=False, w=2.6, rad=2.0):
        x = cx + rad * math.cos(math.radians(ang))
        y = cy + rad * math.sin(math.radians(ang))
        box = s.shapes.add_textbox(Inches(x - w / 2), Inches(y - 0.4), Inches(w), Inches(0.8))
        tf = box.text_frame; tf.word_wrap = True
        lines = text.split("\n") if two_line else [text]
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            rr = p.add_run(); rr.text = ln
            rr.font.size = Pt(size); rr.font.bold = True
            rr.font.color.rgb = color; rr.font.name = FONT
        return box

    def arc(name):
        return s.shapes.add_picture(asset(name),
            Inches(cx - side / 2), Inches(cy - side / 2), Inches(side), Inches(side))
    arc("ring_llm.png")
    arc("ring_agent.png")
    arc("ring_context.png")

    txb(s, "AI", Inches(cx - 0.8), Inches(cy - 0.4), Inches(1.6), Inches(0.8),
        size=32, bold=True, color=FG, align=PP_ALIGN.CENTER)
    polar_label("LLM", 270, 19, ACCENT, w=1.8, rad=1.95)
    polar_label("Agent", 150, 19, SKY, w=2.0, rad=2.35)
    polar_label("Context\nManagement", 30, 16, GOLD, two_line=True, w=2.4, rad=2.35)

    # Right side: the turn. "but," then the human question, revealed on clicks.
    but = txb(s, "but,", Inches(7.4), Inches(2.7), Inches(5.4), Inches(0.9),
              size=44, bold=True, color=MUTED)
    q = txb(s, "where does it leave us, humans?",
            Inches(7.4), Inches(3.7), Inches(5.5), Inches(2.0),
            size=34, bold=True, color=GOLD)

    animate_clicks(s, [[but], [q]])

    add_notes(s,
        "⏱ ~0:30 | Running: 26:00\n\n"
        "Recap the whole picture: this is AI, three pieces working together,\n"
        "LLM, Agent, Context Management.\n"
        "(click) 'but,' ... let it land.\n"
        "(click) where does it leave us, humans? Bridge into the closing part:\n"
        "what stays our job, what our value is.\n\n"
        "Ring visible from start; 'but,' and the question each reveal on a click.")


def s10k_kasparov(prs):
    """Kasparov / centaur story in one slide, revealed in four clicks."""
    s = content_slide(prs, "", "Learning from the past")

    xs = [Inches(0.55), Inches(4.75), Inches(8.95)]
    col_w = Inches(3.83)
    rect(s, Inches(4.6), Inches(1.85), Pt(1), Inches(3.95), fill=LINE_DIM)
    rect(s, Inches(8.8), Inches(1.85), Pt(1), Inches(3.95), fill=LINE_DIM)

    # ── Beat 1: 1997, the machine won (Deep Blue photo) ──
    g1 = []
    g1.append(txb(s, "1997", xs[0], Inches(1.85), col_w, Inches(0.85),
                  size=22, bold=True, color=CORAL, align=PP_ALIGN.CENTER))
    g1.append(add_pic_fit(s, asset("deep_blue.jpg"), xs[0], Inches(2.7), col_w, Inches(2.2)))
    g1.append(txb(s, "Deep Blue beats the world champion.",
                  xs[0], Inches(4.95), col_w, Inches(0.85),
                  size=17, bold=True, color=FG, align=PP_ALIGN.CENTER))

    # ── Beat 2: the reframe (Kasparov photo) ──
    g2 = []
    g2.append(txb(s, "Kasparov's response:\nAdvanced Chess", xs[1], Inches(1.85), col_w, Inches(0.85),
                  size=19, bold=True, color=SKY, align=PP_ALIGN.CENTER))
    g2.append(add_pic_fit(s, asset("kasparov.jpg"), xs[1], Inches(2.7), col_w, Inches(2.2)))
    g2.append(txb(s, "Human + machine, on the same side.",
                  xs[1], Inches(4.95), col_w, Inches(0.85),
                  size=17, bold=True, color=FG, align=PP_ALIGN.CENTER))

    # ── Beat 3: 2005 freestyle (illustration: 2 amateurs + 3 home PCs) ──
    g3 = []
    g3.append(txb(s, "2005", xs[2], Inches(1.85), col_w, Inches(0.85),
                  size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER))
    g3.append(add_pic_fit(s, asset("freestyle_team.png"), xs[2], Inches(2.7), col_w, Inches(2.2)))
    g3.append(txb(s, "2 amateurs + 3 home PCs beat grandmasters and supercomputers.",
                  xs[2], Inches(4.95), col_w, Inches(0.85),
                  size=17, bold=True, color=FG, align=PP_ALIGN.CENTER))

    # ── Beat 4: Kasparov's conclusion (banner) ──
    g4 = []
    g4.append(rect(s, Inches(0.55), Inches(6.0), Inches(12.23), Inches(1.15),
                   fill=BOX_FILL, line=ACCENT, line_w=Pt(1.5)))
    g4.append(txb(s, "Weak human + machine + a better process  >  a strong machine alone.",
                  Inches(0.8), Inches(6.18), Inches(11.7), Inches(0.6),
                  size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER))
    g4.append(txb(s, "Garry Kasparov's conclusion: the edge is the partnership, not raw power.",
                  Inches(0.8), Inches(6.66), Inches(11.7), Inches(0.4),
                  size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER))

    animate_clicks(s, [g1, g2, g3, g4])

    add_notes(s,
        "⏱ ~1:30 | Running: 27:30\n\n"
        "Tell it as a story, four beats:\n"
        "(click) 1997, IBM's Deep Blue beats Kasparov, the reigning world champion.\n"
        "(click) His response wasn't 'are humans obsolete?' It was 'what if we play on the\n"
        "  same side?' He launched Advanced Chess: human + machine.\n"
        "(click) 2005, a freestyle tournament open to anyone. The winners were two amateurs\n"
        "  on three ordinary home PCs, beating grandmasters and far stronger computers.\n"
        "(click) His conclusion: a weak human + machine + a BETTER PROCESS beat a strong\n"
        "  computer alone, and even a strong human with a worse process. The edge was the\n"
        "  partnership and the process, not raw power. Bridge into: where does that leave us?\n\n"
        "Honesty (if asked): Deep Blue was brute-force search, not an LLM; the centaur edge\n"
        "later faded as engines grew overwhelmingly strong. The lesson is the direction, not\n"
        "a fixed seat. Photos: Deep Blue (James the photographer, CC BY 2.0); Kasparov\n"
        "(Grendelkhan, CC BY-SA 4.0), Wikimedia Commons.")


def s10ph_philosophical(prs):
    """Philosophical beat: Nietzsche on meaning, then what AI is and is not (3 reveals)."""
    s = content_slide(prs, "", "Let's get philosophical")

    quote = txb(s,
        "“Only humans give meaning to things. That is why we are called\n"
        "human: the ones who value.”",
        Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.6),
        size=30, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    attrib = txb(s, "Friedrich Nietzsche, 1883 (loosely translated)",
        Inches(0.9), Inches(3.65), Inches(11.5), Inches(0.4),
        size=16, color=MUTED, align=PP_ALIGN.CENTER)

    beat2 = txb(s, "AI is merely a simulation of human thought, but it is not human.",
        Inches(0.7), Inches(4.7), Inches(11.9), Inches(0.7),
        size=24, bold=True, color=SKY, align=PP_ALIGN.CENTER)

    beat3 = txb(s,
        "It cannot give meaning, cannot act with intent, cannot think of new ideas.",
        Inches(0.7), Inches(5.7), Inches(11.9), Inches(1.0),
        size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    animate_clicks(s, [[beat2], [beat3]])

    add_notes(s,
        "⏱ ~1:15 | Running: 28:45\n\n"
        "Read the Nietzsche line slowly (it's a loose translation, not a literal quote):\n"
        "only humans give meaning to things, valuing is what makes us human.\n"
        "(click) AI is a simulation of human thought, an extraordinary one, but it is not\n"
        "  human.\n"
        "(click) So it cannot give meaning, cannot act with intent, cannot originate genuinely\n"
        "  new ideas. That is the gap, and the gap is where our value lives.\n\n"
        "Quote + attribution visible from start; the two AI lines reveal on click.")


def s10e_ideas_spectrum(prs):
    """10e — A world of ideas: tiers from world to you, and what covers each (1 click)."""
    s = content_slide(prs, "", "A world of ideas")

    full = (Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Column centres match generate_ideas_spectrum.py (CX); baseline rule ~3.66" from top.
    CX = [1.893, 4.28, 6.667, 9.053, 11.44]
    coverage = [
        ("the LLM", ACCENT),
        ("coding agents and public skills", SKY),
        ("company agents and skills", SKY),
        ("team skills", GOLD),
        ("only your brain", CORAL),
    ]
    cap_w = 2.32
    # Reveal order: column 0, its caption, column 1, its caption, ... (10 clicks).
    groups = []
    for i, (src, col) in enumerate(coverage):
        col_img = s.shapes.add_picture(asset(f"ideas_col{i}.png"), *full)
        left = Inches(CX[i] - cap_w / 2)
        cov = txb(s, "covered by", left, Inches(3.95), Inches(cap_w), Inches(0.35),
                  size=12, color=MUTED, align=PP_ALIGN.CENTER)
        src_b = txb(s, src, left, Inches(4.32), Inches(cap_w), Inches(1.4),
                    size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        groups.append([col_img])      # the tier (icon + labels + baseline slice)
        groups.append([cov, src_b])   # its "covered by ..." caption

    animate_clicks(s, groups)

    add_notes(s,
        "⏱ ~1:30 | Running: 26:45\n\n"
        "Zoom out. Every idea the AI uses sits somewhere on this line, from what the whole\n"
        "world shares (left) to what only you know (right). Build it up one tier at a time,\n"
        "each tier then its coverage:\n"
        "  (click) World ... (click) covered by the LLM, baked into the weights, for free.\n"
        "  (click) Industry ... (click) coding agents and public skills.\n"
        "  (click) Cato ... (click) company agents and skills.\n"
        "  (click) Team ... (click) team skills.\n"
        "  (click) You ... (click) only your brain. Nothing else has it.\n\n"
        "The further right you go, the more the work is yours to supply. And that work is\n"
        "exactly skills, agents, and context, the things this talk was about.")


def s10d_layer_summary(prs):
    """10d — Three-layer summary + agent diagram"""
    s = content_slide(prs, "PUTTING IT TOGETHER", "The three layers")

    layers = [("Context Management", MUTED), ("Agent", FG), ("LLM", ACCENT)]
    lft, lft_w = Inches(0.4), Inches(5.2)
    layer_h, gap = Inches(1.3), Inches(0.08)
    y_bottom = Inches(6.1)
    agent_box_top = None
    for i, (name, col) in enumerate(layers):
        y = y_bottom - i * (layer_h + gap)
        if name == "Agent":
            agent_box_top = y
        rect(s, lft, y, lft_w, layer_h, fill=BOX_FILL, line=col,
             line_w=Pt(1.5) if col != MUTED else Pt(0.5))
        txb(s, name, lft + Inches(0.2), y + Inches(0.38),
            lft_w - Inches(0.4), Inches(0.6),
            size=22, bold=True, color=col, align=PP_ALIGN.CENTER)

    exp_y = agent_box_top + layer_h / 2 - Inches(0.3)
    txb(s, "«", lft + lft_w + Inches(0.05), exp_y, Inches(0.7), Inches(0.6),
        size=28, bold=True, color=FG, align=PP_ALIGN.CENTER)

    img_l = lft + lft_w + Inches(0.65)
    s.shapes.add_picture(asset("agent_diagram_4.png"),
                         img_l, Inches(1.9), Inches(13.33) - img_l - Inches(0.2), Inches(5.3))
    add_notes(s,
        "⏱ 0:30 | Running: 26:00\n\n"
        "Quick callback. LLM at the bottom (autocomplete at scale), Agent in the middle\n"
        "(LLM + tools + loop), Context Management on top (what the LLM sees). These three\n"
        "describe every AI product you'll use or build. Now — what to do with them.")


def s11_tips(prs):
    """11 — Takeaways"""
    s = content_slide(prs, "TAKEAWAYS", "Using AI well")
    tips = [
        ('"Give me the right context\nand I shall move the world"',
         "AI can do almost anything —\nif you manage its context correctly."),
        ("Better agent = less manual context",
         "The best agents need the least hand-holding.\nThat's what you're optimising for."),
        ("AI disappoints you?",
         "Ask yourself: what did I forget to give it?"),
    ]
    y = Inches(2.0)
    for header, body in tips:
        txb(s, header, Inches(0.55), y, Inches(5.8), Inches(0.9),
            size=20, bold=True, color=ACCENT, italic=True)
        txb(s, body, Inches(6.5), y + Inches(0.05), Inches(6.5), Inches(0.85),
            size=18, color=FG)
        y += Inches(1.6)
    add_notes(s,
        "⏱ 3:00 | Running: 26:15\n\n"
        "Take your time — this is the payoff.\n"
        "1. Archimedes: 'Give me a lever long enough...' — give the LLM the right\n"
        "context and it can do almost anything. The limit is the context, not the model.\n"
        "2. Better agent = less context management: a great agent figures it out from\n"
        "minimal context. That's what you're building towards.\n"
        "3. If AI disappoints — ask what you forgot to give it. 9/10 it's the context.")


def s10f_ai_helps(prs):
    """AI can help you: three ways, revealed one per click."""
    s = content_slide(prs, "", "AI can help you")

    bullet_shapes = []
    y = Inches(2.4)
    for text in ["By doing things you are bad at",
                 "By accelerating things you are good at",
                 "By helping you learn new things and ideas"]:
        tb = s.shapes.add_textbox(Inches(1.2), y, Inches(11.0), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        m = p.add_run(); m.text = "▸   "
        m.font.size = Pt(28); m.font.bold = True; m.font.color.rgb = ACCENT; m.font.name = FONT
        r = p.add_run(); r.text = text
        r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = FG; r.font.name = FONT
        bullet_shapes.append(tb)
        y += Inches(1.25)

    animate_clicks(s, [[b] for b in bullet_shapes])

    add_notes(s,
        "⏱ ~1:00 | Running: 27:30\n\n"
        "The optimistic turn. AI helps you three ways:\n"
        "(click) it does the things you're bad at,\n"
        "(click) it accelerates the things you're good at,\n"
        "(click) and it helps you learn new things and ideas.\n\n"
        "That last one sets up the first closing quote (Mark Cuban) on the next slide.")


def s10g_closing_quotes(prs):
    """Closing finale: four parting quotes as a 2x2 grid, one per click."""
    s = content_slide(prs, "", "A few parting thoughts")

    # (quote, attribution, border colour), in reveal/reading order TL, TR, BL, BR.
    quotes = [
        ("“We can know more than we can tell.”",
         "Michael Polanyi, The Tacit Dimension, 1966", SKY),
        ("“When anyone can make anything, knowing what is worth making is the rare skill.”",
         "echoing Paul Graham, Taste for Makers, 2002", GOLD),
        ("“There are 2 types of LLM users: those who use it to learn everything, "
         "and those who use it so they don't have to learn anything.”",
         "Mark Cuban, 2026", ACCENT),
        ("“A computer can never be held accountable. That's your job as the human in the loop.”",
         "Simon Willison, 2025", CORAL),
    ]

    cw, ch = Inches(5.94), Inches(2.45)
    xs = [Inches(0.55), Inches(6.84)]
    ys = [Inches(1.9), Inches(4.7)]
    pad = Inches(0.3)
    groups = []
    for i, (quote, who, col) in enumerate(quotes):
        x = xs[i % 2]
        yv = ys[i // 2]
        box = rect(s, x, yv, cw, ch, fill=BOX_FILL, line=col, line_w=Pt(1.5))
        q = txb(s, quote, x + pad, yv + Inches(0.22), cw - 2 * pad, Inches(1.55),
                size=17, italic=True, color=FG)
        a = txb(s, who, x + pad, yv + ch - Inches(0.55), cw - 2 * pad, Inches(0.4),
                size=13, bold=True, color=col)
        groups.append([box, q, a])

    animate_clicks(s, groups)

    add_notes(s,
        "⏱ ~1:30 | Running: 29:00\n\n"
        "Four voices to leave them with, one per click:\n"
        "(click) Polanyi: we know more than we can tell, your tacit knowledge is yours\n"
        "  and can't be handed over. (The Tacit Dimension, 1966)\n"
        "(click) Graham: when anyone can make anything, taste, knowing what's worth\n"
        "  making, is the rare skill. (echoing Taste for Makers, 2002)\n"
        "(click) Cuban: use AI to learn everything, not to avoid learning. (2026)\n"
        "(click) Willison: a computer can never be held accountable, that's your job as\n"
        "  the human in the loop. (2025; the first line is a 1979 IBM training slide,\n"
        "  Willison's contribution is the 'human in the loop' framing.)\n\n"
        "End on Willison: the responsibility stays on us. Then: Thank you.")


def s12_thanks(prs):
    """12 — Thank you (Thank you slide Black)"""
    s = add(prs, "Thank you slide Black")
    set_ph(s, 0, "Thank you", color=FG)

    # The cute geek stickers from the title slide, scattered around the "Thank you"
    # (kept clear of the centre title band, x 3.3-10.0 / y 3.3-4.75).
    geeks = [
        ("geek_claude_t.png",   0.70, 0.65, 1.75, -14),
        ("geek_copilot_t.png", 10.85, 0.55, 1.50,  15),
        ("geek_cursot_t.png",   0.45, 3.05, 1.35,   9),
        ("geek_claude_t.png",  11.45, 3.15, 1.35,  -9),
        ("geek_cursot_t.png",   1.35, 5.35, 1.60,  11),
        ("geek_copilot_t.png", 10.40, 5.45, 1.45, -12),
    ]
    for name, l, t, wd, rot in geeks:
        p = asset(name)
        if os.path.exists(p):
            pic = s.shapes.add_picture(p, Inches(l), Inches(t), width=Inches(wd))
            pic.rotation = rot

    add_notes(s,
        "⏱ 3:45 | Running: 30:00\n\n"
        "Close warm. The geek stickers from the opening come back, scattered, to bookend\n"
        "the talk. Optional callback before Q&A: 'You now know the parts. Go fire the\n"
        "weapon.' Budget ~3-4 min for Q&A. Expect: when to use an agent vs a plain LLM\n"
        "call (real-world data / multi-step), best model (depends; concepts are the same),\n"
        "how we use this at Cato (good opening for a follow-up session).")


# ═══════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    s01_title(prs)
    s02_intro_me(prs)
    s03_geek_origin(prs)
    s04a_done_before_grid(prs)       # variant A1 — pick grid or timeline, remove the other
    s04b_done_before_timeline(prs)   # variant A2
    s05_ai_different(prs)
    s07_breakdown_rings(prs)         # replaces pin-reveal + the layer-stack roadmap
    # ── Piece 1 · LLM ────────────────────────────────────────────────────────
    s08_llm_opener(prs)              # LLM = Large Language Model; two phases
    s08t_training(prs)               # Phase 1: Training (llm_training.png)
    s08op_operation(prs)             # Phase 2: black-box autocomplete (Cato joke)
    s08p1_trivial_logic(prs)         # property: trivial + logic
    s08p2_knowledge_frozen(prs)      # property: knowledge + frozen in time
    s08p3_statistical_predicts(prs)  # property: statistical + predicts, not lookup
    s08p4_subjective_oneletter(prs)  # property: subjective + one letter
    s08p5_every_letter(prs)          # property: one letter changes everything (cat vs car)
    s08p6_roleplay(prs)              # property: role-play / persona, bridges into Agents
    # ── Piece 2 · AGENT ──────────────────────────────────────────────────────
    s09a0_context_limit(prs)         # LLM limited by input size + growth chart on click
    s09a0c_too_big(prs)              # ...but some things still do not fit (GoT tokens)
    s09mr_mapreduce(prs)             # V2: many LLMs + plain code (map-reduce tree)
    s09loop_notools(prs)             # V3: feed output back in, the exit-point question
    s09loop_tools(prs)               # V4: add tools -> an agent (full loop diagram)
    s09b_developer_wrote(prs)        # walk-through opens: the user's prompt
    s09b1_call0(prs)                 # CALL 0: the assembled [SYSTEM]+[HUMAN] prompt
    s09c_call1(prs)                  # CALL 1: LLM decided (web_search)
    s09d_call2(prs)                  # CALL 2: + legend
    s09e_call3(prs)
    s09f_call4(prs)
    s09cloudlocal(prs)               # cloud vs local agents (brain usually cloud)
    # ── Piece 3 · CONTEXT MANAGEMENT ─────────────────────────────────────────
    s10a_context_why(prs)            # why manage context (cost AND quality) + what to watch
    s10a2_context_moves(prs)         # the three moves: start fresh / compress / offload
    s10b_skills_filesystem(prs)      # worked example: skills (offload + load on demand)
    s10c_ecosystem(prs)
    s10r_so_this_is_ai(prs)          # recap ring (small, left) -> "but, where does it leave us?"
    s10k_kasparov(prs)               # Kasparov / centaur story (4 clicks), ends on his conclusion
    s10ph_philosophical(prs)         # Nietzsche on meaning -> AI is not human, cannot give meaning
    s10e_ideas_spectrum(prs)         # zoom out: world → you, and how each tier reaches the model
    s10f_ai_helps(prs)               # the optimistic turn: three ways AI helps you
    s10g_closing_quotes(prs)         # finale: four parting quotes (Cuban, Polanyi, Graham, Willison)
    s12_thanks(prs)

    out = os.path.join(HERE, "understanding_ai_for_geeks.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
    return out


if __name__ == "__main__":
    build()
